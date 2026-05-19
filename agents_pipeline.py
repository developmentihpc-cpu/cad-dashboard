"""
CAD Development Intelligence Agents
=====================================
Paste this file into your Claude Code project alongside digest_agent.py.

Usage (trigger by email):
  python agents_pipeline.py                        # manual run, no email
  python agents_pipeline.py --email                # run + send PPT by email
  python agents_pipeline.py --watch                # poll inbox every 5 min

Trigger email format (send to GMAIL_SENDER):
  Subject: RUN: <country> / <sector> / <workflow>

  Workflow values:
    needs       -> Country Needs Assessment (9 agents, 13 slides)
    proposal    -> Development Proposal     (11 agents, 16 slides)
    external    -> External Assessment      (3 agents, scorecard)

  Examples:
    RUN: Senegal / Education & Health / needs
    RUN: Ethiopia / WASH & Climate / proposal
    RUN: Somalia / All sectors / external

Setup:
  pip install anthropic
  Add ANTHROPIC_API_KEY to your .env file
"""

import json, os, re, sys, smtplib, time, base64, subprocess, tempfile
from datetime import datetime
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.base import MIMEBase
from email import encoders
from pathlib import Path
import anthropic
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Config ────────────────────────────────────────────────────────────────────

BASE_DIR   = Path(__file__).parent
ENV_FILE   = BASE_DIR / ".env"
STATE_FILE = BASE_DIR / "pipeline-state.json"
OUT_DIR    = BASE_DIR / "agent_outputs"
OUT_DIR.mkdir(exist_ok=True)

POLL_INTERVAL_SECONDS = 300   # 5 minutes
TRIGGER_PREFIX        = "RUN:"

# Template palette
T = {
    "navy":       "1D252C",
    "gold":       "AD833B",
    "gold_light": "C7A877",
    "blue":       "2F586E",
    "steel":      "678CA5",
    "steel_lt":   "CBDCE6",
    "deep_red":   "79242F",
    "navy2":      "333F64",
    "gray":       "4A4F54",
    "bg":         "F5F7F9",
    "white":      "FFFFFF",
    "black":      "101820",
    "pos":        "00BC8B",
    "semi_pos":   "9CBB5D",
    "neutral":    "F39B26",
    "semi_neg":   "BA492F",
    "neg":        "9B120B",
    "text":       "1D252C",
    "muted":      "4A4F54",
    "border":     "E5E8EF",
}

STEP_NAMES = ["Intake","Country Analysis","Program Dev",
              "Stakeholder Eng.","Leadership Review","Approved"]

# ── Env loader ─────────────────────────────────────────────────────────────────

def load_env():
    env = {}
    if ENV_FILE.exists():
        for line in ENV_FILE.read_text(encoding="utf-8").splitlines():
            line = line.strip()
            if line and not line.startswith("#") and "=" in line:
                k, v = line.split("=", 1)
                env[k.strip()] = v.strip()
    for key in ("GMAIL_SENDER","GMAIL_APP_PASSWORD","DIGEST_RECIPIENT","ANTHROPIC_API_KEY"):
        if key not in env and key in os.environ:
            env[key] = os.environ[key]
    return env

ENV = load_env()

# ── Anthropic client ──────────────────────────────────────────────────────────

def get_client():
    api_key = ENV.get("ANTHROPIC_API_KEY","")
    if not api_key or api_key.startswith("sk-ant-xxx"):
        print("[ERROR] ANTHROPIC_API_KEY not set in .env")
        print("  Add:  ANTHROPIC_API_KEY=sk-ant-...")
        sys.exit(1)
    return anthropic.Anthropic(api_key=api_key)

# ══════════════════════════════════════════════════════════════════════════════
# AGENT DEFINITIONS
# ══════════════════════════════════════════════════════════════════════════════

NEEDS_AGENTS = [
    "router", "needs_orch", "fact_checker",
    "analyst", "strategist", "content", "editor"
]

PROPOSAL_AGENTS = [
    "router", "prop_orch", "fact_checker",
    "analyst", "strategist", "proj_design", "costing",
    "content", "editor"
]

EXTERNAL_AGENTS = [
    "router", "ext_assessor", "editor"
]

WORKFLOW_AGENTS = {
    "needs":    NEEDS_AGENTS,
    "proposal": PROPOSAL_AGENTS,
    "external": EXTERNAL_AGENTS,
}

WORKFLOW_SLIDES = {
    "needs": 13,
    "proposal": 16,
    "external": None,
}

def load_country_context(country):
    """
    Load exported country-context.json from the project folder if available.
    Returns formatted string summary for injection into agent prompts.
    """
    # Look for country-specific context file
    safe_name = country.replace(' ', '_')
    candidates = [
        BASE_DIR / f"{safe_name}_context.json",
        BASE_DIR / "agent_outputs" / f"{safe_name}_context.json",
    ]
    for path in candidates:
        if path.exists():
            try:
                data = json.loads(path.read_text(encoding="utf-8"))
                indicators = data.get("indicators", {})
                pipeline   = data.get("pipeline", [])
                programs   = data.get("programs", {})
                info       = data.get("countryInfo", {})
                quality    = data.get("dataQuality", {})

                # Format indicators
                ind_lines = []
                label_map = {
                    "NY.GDP.PCAP.PP.CD": "GDP/capita PPP (USD)",
                    "SI.POV.DDAY":       "Extreme poverty rate (%)",
                    "SE.ADT.LITR.ZS":    "Adult literacy rate (%)",
                    "SH.DYN.MORT":       "Child mortality (/1,000)",
                    "EG.ELC.ACCS.ZS":    "Electricity access (%)",
                    "NY.GDP.MKTP.KD.ZG": "GDP growth rate (%)",
                    "FP.CPI.TOTL.ZG":    "Inflation rate (%)",
                    "SP.POP.TOTL":       "Total population",
                    "SH.H2O.SMDW.ZS":   "Safe water access (%)",
                    "SN.ITK.DEFC.ZS":   "Undernourishment (%)",
                    "SI.POV.GINI":       "Gini index",
                    "SH.XPD.CHEX.GD.ZS":"Health expenditure (% GDP)",
                    "SE.XPD.TOTL.GD.ZS":"Education expenditure (% GDP)",
                }
                for ind_id, val in indicators.items():
                    label = label_map.get(ind_id, ind_id)
                    ind_lines.append(f"  {label}: {val['value']} ({val.get('year','')})")

                # Format pipeline
                pip_lines = []
                for p in pipeline:
                    scores = ""
                    if p.get("evalScores"):
                        scores = " | Eval: " + ", ".join(
                            f"{s['area']} {s['rating']}" for s in p["evalScores"][:3])
                    pip_lines.append(
                        f"  - {p['title']} ({p['sector']}, {p['path']}, "
                        f"{p['status']}, {p['cost']}, {p.get('beneficiaries','—')} ben.){scores}"
                    )

                # Format programs
                prog_lines = []
                for ptype in ["existing", "proposed", "evaluation"]:
                    for pg in programs.get(ptype, [])[:3]:
                        prog_lines.append(f"  - [{ptype}] {pg.get('nameEn','—')} ({pg.get('sector','—')}, {pg.get('cost','—')})")

                parts = [f"=== DASHBOARD CONTEXT FOR {country.upper()} ===",
                         f"Exported: {data.get('generatedAt','')[:10]}",
                         f"Capital: {info.get('capital','—')} | Currency: {info.get('currency','—')} | Languages: {info.get('languages','—')}",
                         f"Data quality: {quality.get('indicatorCount',0)} indicators, {quality.get('pipelineProjects',0)} pipeline projects\n"]

                if ind_lines:
                    parts.append("WORLD BANK INDICATORS (verified, live data):")
                    parts.extend(ind_lines)
                if pip_lines:
                    parts.append(f"\nCAD PIPELINE PROJECTS ({len(pipeline)}):")
                    parts.extend(pip_lines)
                if prog_lines:
                    parts.append(f"\nCOUNTRY PROGRAMS:")
                    parts.extend(prog_lines)

                parts.append("\n=== USE THE ABOVE AS YOUR PRIMARY DATA SOURCE ===")
                print(f"[OK] Loaded dashboard context for {country} ({quality.get('indicatorCount',0)} indicators)")
                return "\n".join(parts)
            except Exception as e:
                print(f"[WARN] Could not load context for {country}: {e}")
    return ""   # No context file found — agents will use their own knowledge


def build_system_prompt(agent_id, ctx):
    country  = ctx["country"]
    sector   = ctx["sector"]
    audience = ctx.get("audience", "Senior Leadership & Donor Partners")
    budget   = ctx.get("budget", "TBD")
    horizon  = ctx.get("horizon", "3 years")
    workflow = ctx["workflow"]
    slides   = WORKFLOW_SLIDES.get(workflow)
    slide_str = f"{slides}-slide " if slides else ""

    # Load dashboard-exported context (replaces research agent)
    dashboard_ctx = load_country_context(country)
    context_note  = (
        "\n\nDASHBOARD DATA AVAILABLE: You have been provided with verified World Bank "
        "indicators and CAD pipeline data exported directly from the dashboard. "
        "Use this as your primary data source. Do not guess or hallucinate figures."
        if dashboard_ctx else
        "\n\nNO DASHBOARD EXPORT FOUND: Use your best knowledge of this country. "
        "Clearly flag any figures as estimates and recommend verification."
    )
    full_ctx = (dashboard_ctx + context_note) if dashboard_ctx else context_note

    prompts = {

        "router": (
            f"You are the Router in the Country Assessment Dashboard (CAD) agent system. "
            f"Classify this request: workflow={workflow}, country={country}, sector={sector}, "
            f"audience={audience}. "
            f"Confirm routing to the correct workflow "
            f"({'country_needs_report' if workflow=='needs' else 'country_development_proposal' if workflow=='proposal' else 'external_project_criteria_assessment'}). "
            f"State assumptions, confirm the pipeline will proceed. Be concise (3 paragraphs)."
        ),

        "needs_orch": (
            f"You are the Country Needs Report Orchestrator in the CAD system. "
            f"Define the execution plan for a {slide_str}needs assessment of {country}, "
            f"sector: {sector}, audience: {audience}. "
            f"Output: confirmed scope, required data sources (World Bank, UNDP, WHO, UNESCO, "
            f"FAO, WFP, national development plan), 13-slide structure, quality gates, "
            f"key assumptions, agent task sequence. No project design.{full_ctx}"
        ),

        "prop_orch": (
            f"You are the Country Development Proposal Orchestrator in the CAD system. "
            f"Define the execution plan for a {slide_str}development proposal for {country}, "
            f"sector: {sector}, budget: {budget}, horizon: {horizon}, audience: {audience}. "
            f"Output: scope, data sources, 16-slide structure, quality gates, donor alignment "
            f"notes, agent task sequence.{full_ctx}"
        ),

        "fact_checker": (
            f"You are the Fact Checker in the CAD system. "
            f"Your job is narrow and specific: review the dashboard-exported data for {country} "
            f"and flag any issues before the analysis begins.\n\n"
            f"{full_ctx}\n\n"
            f"Output exactly four sections:\n"
            f"1. DATA COMPLETENESS — which key indicators are missing or outdated (>3 years old). "
            f"Flag if poverty, GDP, literacy, or mortality data is absent.\n"
            f"2. KNOWN CONTEXT GAPS — critical facts NOT in the dashboard data: conflict status, "
            f"political context, major donor programmes, recent crises. Keep to 3-5 bullet points.\n"
            f"3. DATA QUALITY VERDICT — one of: SUFFICIENT / SUFFICIENT WITH GAPS / INSUFFICIENT. "
            f"State what the analyst should treat as estimates.\n"
            f"4. APPROVED DATA SUMMARY — a clean, concise list of the verified figures the "
            f"downstream agents should use. Do not add figures you are not confident about."
        ),

        "analyst": (
            f"You are the Needs and Benchmarking Analyst in the CAD system. "
            f"Using the verified data below, identify key development gaps in {sector} for {country}. "
            f"Benchmark vs. regional peers and SDG targets. Map geographic disparities "
            f"and vulnerable groups. Identify root causes - not just symptoms. "
            f"Produce severity and urgency rankings. "
            f"Link gaps to potential intervention areas - do NOT design projects or budgets.\n\n"
            f"{full_ctx}"
        ),

        "strategist": (
            f"You are the Strategic Prioritization Specialist in the CAD system. "
            f"For {country} ({sector}{',' + budget if budget != 'TBD' else ''}): "
            f"apply a transparent scoring matrix (severity, feasibility, impact, "
            f"donor alignment, sustainability, cost-effectiveness, SDG relevance). "
            f"Output: scoring matrix, top 3 priority areas with full rationale, "
            f"excluded areas with reasons, trade-offs, strategic narrative, "
            f"donor and national plan alignment. Do not recommend everything.\n\n"
            f"{full_ctx}"
        ),

        "proj_design": (
            f"You are the Project Design Specialist in the CAD system. "
            f"Design 3 evidence-based project concepts for {country} in {sector}. "
            f"For each: title, problem statement with evidence links, objective, "
            f"target geography and beneficiaries, components, key activities, "
            f"theory of change, implementation model, potential partners, "
            f"sustainability approach, expected outputs/outcomes, phasing, "
            f"assumptions, risks, SDG links. "
            f"Projects must be practical and fundable. Total budget: {budget}."
        ),

        "costing": (
            f"You are the Costing, Results and Risk Specialist in the CAD system. "
            f"For 3 projects in {country}: produce indicative budget breakdown by project "
            f"and component (total ~{budget}), cost assumptions, cost-per-beneficiary, "
            f"results framework, KPI table with baselines and targets, logframe elements, "
            f"full risk matrix with mitigation measures, safeguards notes, "
            f"sustainability risks. Flag assumptions and where local validation is needed."
        ),

        "content": (
            f"You are the PowerPoint Content Builder in the CAD system. "
            f"Build a {slide_str}content package for {country} ({workflow}). "
            + (
                "Structure: 1.Title, 2.Exec Summary, 3.Country Context, 4.Socioeconomic Snapshot, "
                "5.National Priorities, 6.Sector Needs, 7.SDG Benchmarking, 8.Vulnerable Groups, "
                "9.Root Causes, 10.Priority Ranking, 11.Strategic Implications, "
                "12.Evidence Gaps, 13.Conclusion."
                if workflow == "needs" else
                "Structure: 1.Title, 2.Exec Summary, 3.Country Context, 4.Key Challenges, "
                "5.Needs Snapshot, 6.Strategic Priorities, 7.Program Overview, "
                "8.Project 1, 9.Project 2, 10.Project 3, 11.Budget Summary, "
                "12.Results & KPIs, 13.Implementation Timeline, 14.Risk Matrix, "
                "15.Sustainability, 16.Next Steps."
            )
            + f" For each slide: number, title, key message, 3-4 bullets, suggested visual. "
            f"Keep slides concise. No unsupported claims."
            + (" No project proposals." if workflow == "needs" else "")
        ),

        "ppt_designer": (
            f"You are the PowerPoint Designer in the CAD system. "
            f"Produce a detailed design blueprint for the {country} {workflow} deck. "
            f"Template palette: Navy #1D252C, Gold #AD833B, Deep blue #2F586E, "
            f"Steel blue #678CA5, White #FFFFFF. Font: Montserrat. "
            f"Cover: white background, light gray left panel (35%), geometric circle motif. "
            f"Content slides: white background, dark navy header bar, gold 3px left accent stripe. "
            f"Charts use blue primary bars, severity badges in red/amber/green."
        ),

        "editor": (
            f"You are the Final Quality Editor in the CAD system. "
            f"Run the complete quality checklist for the {country} {workflow} output. "
            f"Check: correct workflow used, narrative coherent, evidence cited, "
            f"benchmarks clear, priorities justified"
            + (", projects practical and fundable, budgets realistic and caveated, "
               "KPIs match activities" if workflow == "proposal" else
               ", no project proposals included")
            + f", figures consistent, no unsupported claims, leadership/donor-ready. "
            f"Issue final verdict: READY, READY WITH CAVEATS, or NOT READY. "
            f"List remaining issues concisely."
        ),

        "ext_assessor": (
            f"You are the External Project Criteria Assessor in the CAD system. "
            f"Assess {ctx.get('num_projects', 3)} submitted external projects for "
            f"{country} in {sector}. "
            f"Apply scoring framework: strategic alignment, country priority alignment, "
            f"sector relevance, evidence quality, development impact, beneficiary clarity, "
            f"feasibility, budget reasonableness, implementation capacity, risk level, "
            f"sustainability, M&E logic, safeguards. "
            f"For each project produce: summary, score /100, alignment assessment, "
            f"red flags, missing information, recommendation (approve/revise/reject/hold). "
            f"Then produce ranked list with overall recommendations and next steps."
        ),
    }

    return prompts.get(agent_id, f"You are the {agent_id} agent for {country}.")


# ══════════════════════════════════════════════════════════════════════════════
# PIPELINE RUNNER
# ══════════════════════════════════════════════════════════════════════════════

def run_pipeline(ctx, on_agent_start=None, on_agent_done=None):
    workflow = ctx.get("workflow", "needs")
    agents   = WORKFLOW_AGENTS.get(workflow, NEEDS_AGENTS)
    client   = get_client()
    outputs  = {}

    print(f"\n{'='*56}")
    print(f"  CAD Agent Pipeline -- {workflow.upper()}")
    print(f"  Country: {ctx['country']}  |  Sector: {ctx['sector']}")
    print(f"  {len(agents)} agents  |  {datetime.now().strftime('%Y-%m-%d %H:%M')}")
    print(f"{'='*56}")

    for i, agent_id in enumerate(agents):
        step  = i + 1
        total = len(agents)

        if on_agent_start:
            on_agent_start(agent_id, step, total)

        print(f"\n[{step}/{total}] {agent_id} ...", end=" ", flush=True)

        prev_ctx = "\n\n---\n\n".join(
            f"[{k.upper()} OUTPUT]\n{v}"
            for k, v in outputs.items()
        )
        user_msg = (
            f"Previous agent outputs for context:\n\n{prev_ctx}\n\n"
            f"---\nNow complete your role as {agent_id} for {ctx['country']}."
            if prev_ctx else
            f"Complete your role as {agent_id} for the {workflow} "
            f"workflow -- country: {ctx['country']}."
        )

        full_text = ""
        with client.messages.stream(
            model="claude-opus-4-5",
            max_tokens=1500,
            system=build_system_prompt(agent_id, ctx),
            messages=[{"role": "user", "content": user_msg}],
        ) as stream:
            for text in stream.text_stream:
                full_text += text
                print(".", end="", flush=True)

        outputs[agent_id] = full_text
        print(f" done ({len(full_text)} chars)")

        out_path = OUT_DIR / f"{ctx['country'].lower().replace(' ','_')}_{workflow}_{agent_id}.txt"
        out_path.write_text(full_text, encoding="utf-8")

        if on_agent_done:
            on_agent_done(agent_id, step, total, full_text)

    print(f"\n[OK] Pipeline complete -- {len(agents)} agents ran")
    return outputs


# ══════════════════════════════════════════════════════════════════════════════
# PPT GENERATOR (python-pptx, no Node.js required)
# ══════════════════════════════════════════════════════════════════════════════

def rgb(hex_str):
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def add_text_box(slide, text, left, top, width, height,
                 font_size=11, bold=False, color="#1D252C",
                 align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return txBox

def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_color)
    if line_color:
        shape.line.color.rgb = rgb(line_color)
    else:
        shape.line.fill.background()
    return shape

def extract_bullets(text, max_bullets=5, max_chars=160):
    """Extract clean bullet points from agent output text."""
    lines = [l.strip() for l in text.splitlines() if l.strip()]
    bullets = []
    for line in lines:
        # Skip headers, short lines, markdown artifacts
        if len(line) < 20:
            continue
        if line.startswith("#") or line.startswith("---"):
            continue
        # Clean markdown
        line = re.sub(r'\*\*(.+?)\*\*', r'\1', line)
        line = re.sub(r'^\s*[-*•]\s*', '', line)
        line = re.sub(r'^\d+\.\s*', '', line)
        if len(line) > max_chars:
            line = line[:max_chars] + "..."
        if len(line) >= 20:
            bullets.append(line)
        if len(bullets) >= max_bullets:
            break
    return bullets if bullets else ["See full agent output in agent_outputs/ folder."]

def add_content_slide(prs, slide_num, title, tag, bullets, verdict_color=None):
    """Standard content slide: white bg, navy header, gold left stripe."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(T["white"])

    # Navy header bar
    add_rect(slide, 0, 0, 10, 0.55, T["navy"])
    # Gold left accent stripe
    add_rect(slide, 0, 0.55, 0.04, 4.82, T["gold"])

    # Tag (small caps above title)
    if tag:
        add_text_box(slide, tag, 0.14, 0.04, 6, 0.25,
                     font_size=6.5, bold=True, color=T["gold_light"])
    # Title
    add_text_box(slide, title, 0.14, 0.27, 9.5, 0.3,
                 font_size=14, bold=True, color=T["white"])

    # Bullet content area
    y = 0.72
    for i, bullet in enumerate(bullets[:6]):
        # Bullet row background (alternating)
        bg_color = T["bg"] if i % 2 == 0 else T["white"]
        add_rect(slide, 0.14, y, 9.72, 0.58, bg_color)
        # Gold dot
        add_rect(slide, 0.18, y + 0.22, 0.06, 0.06, T["gold"])
        # Bullet text
        add_text_box(slide, bullet, 0.32, y + 0.08, 9.4, 0.44,
                     font_size=10, color=T["text"])
        y += 0.62

    # Footer
    add_rect(slide, 0, 5.3, 10, 0.33, T["bg"])
    add_text_box(slide, f"Office of Development Affairs  ·  Country Assessment Dashboard  ·  {datetime.now().strftime('%B %Y')}",
                 0.14, 5.3, 8.5, 0.33, font_size=6.5, color=T["muted"])
    add_text_box(slide, str(slide_num).zfill(2),
                 9.2, 5.3, 0.65, 0.33,
                 font_size=7, bold=True, color=T["gold"], align=PP_ALIGN.RIGHT)
    return slide

def generate_pptx(ctx, agent_outputs, verdict):
    """Generate the 12-slide country brief via the canonical builder.

    Path 3 of the unified PPT architecture — see country_ppt_builder.py.
    Loads dashboard-exported JSON when present, merges agent outputs into
    the builder context, then delegates to country_ppt_builder.build().

    Returns the output file path.
    """
    from agent_to_builder_ctx import build_ctx_from_agent_run
    from country_ppt_builder import build as build_ppt

    country  = ctx["country"]
    workflow = ctx["workflow"]
    fname    = f"{country.replace(' ','_')}_{workflow}_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx"
    out_path = OUT_DIR / fname

    builder_ctx = build_ctx_from_agent_run(ctx, agent_outputs, verdict)
    build_ppt(builder_ctx, str(out_path))
    size_kb = out_path.stat().st_size // 1024 if out_path.exists() else 0
    print(f"[OK] Generated {fname} ({size_kb}KB) via country_ppt_builder")
    return out_path


def _legacy_generate_pptx(ctx, agent_outputs, verdict):  # pragma: no cover
    """Original agent-output-per-slide layout. Kept for fallback debugging only."""
    country  = ctx["country"]
    workflow = ctx["workflow"]
    sector   = ctx["sector"]
    audience = ctx.get("audience", "Senior Leadership & Donor Partners")
    date_str = datetime.now().strftime("%B %Y")
    fname    = f"{country.replace(' ','_')}_{workflow}_{datetime.now().strftime('%Y%m%d_%H%M')}_legacy.pptx"
    out_path = OUT_DIR / fname

    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)

    # ── SLIDE 1: Cover ────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(T["white"])

    # Left panel
    add_rect(slide, 0, 0, 3.2, 5.625, T["bg"])
    add_rect(slide, 0, 0, 0.04, 5.625, T["gold"])

    # Circle motifs
    for r_val, lc in [(1.8, T["gold"]), (1.3, T["steel"]), (0.85, T["steel_lt"])]:
        cx = 1.6 - r_val; cy = 2.4 - r_val
        shape = slide.shapes.add_shape(9, Inches(cx), Inches(cy),
                                        Inches(r_val*2), Inches(r_val*2))
        shape.fill.background()
        shape.line.color.rgb = rgb(lc)
        shape.line.width = Pt(0.8)

    # Title block
    add_rect(slide, 3.4, 0.65, 0.04, 0.55, T["gold"])
    add_text_box(slide, country.upper(), 3.55, 0.65, 5.5, 0.25,
                 font_size=7.5, bold=True, color=T["gold"])
    add_text_box(slide, date_str, 3.55, 0.9, 5.5, 0.22,
                 font_size=8, color=T["steel"])

    wf_titles = {"needs":"Country Needs Assessment",
                 "proposal":"Development Proposal",
                 "external":"External Project Assessment"}
    add_text_box(slide, wf_titles.get(workflow, "Assessment"),
                 3.55, 1.2, 5.8, 0.65,
                 font_size=26, bold=True, color=T["text"])
    add_text_box(slide, sector, 3.55, 2.0, 5.8, 0.35,
                 font_size=11, bold=True, color=T["blue"])
    add_rect(slide, 3.55, 2.5, 0.65, 0.04, T["gold"])
    add_text_box(slide, f"Prepared for: {audience}", 3.55, 2.62, 5.8, 0.22,
                 font_size=8, color=T["muted"])

    agents_used = WORKFLOW_AGENTS.get(workflow, NEEDS_AGENTS)
    add_text_box(slide, f"Workflow: {workflow.upper()}  ·  {len(agents_used)} specialist agents  ·  Confidential",
                 3.55, 2.86, 5.8, 0.22, font_size=8, color=T["muted"])

    # Verdict badge
    vc = {"READY": T["pos"], "READY WITH CAVEATS": T["neutral"], "NOT READY": T["neg"]}
    add_rect(slide, 3.55, 3.25, 1.6, 0.3, vc.get(verdict, T["pos"]))
    add_text_box(slide, f"✓  {verdict}", 3.55, 3.25, 1.6, 0.3,
                 font_size=7.5, bold=True, color=T["white"], align=PP_ALIGN.CENTER)

    # Black footer
    add_rect(slide, 0, 5.3, 10, 0.325, T["black"])
    add_text_box(slide, f"CONFIDENTIAL  ·  OFFICE OF DEVELOPMENT AFFAIRS  ·  {date_str.upper()}",
                 0.14, 5.3, 9.7, 0.325, font_size=6.5, bold=True, color=T["gold"])

    # ── SLIDES 2+: One slide per key agent output ─────────────────────────────
    slide_defs = {
        "needs": [
            ("research",    "Country Evidence Pack",        "RESEARCH"),
            ("verifier",    "Data Verification",            "VERIFICATION"),
            ("analyst",     "Needs & Benchmarking Analysis","ANALYSIS"),
            ("strategist",  "Strategic Priority Ranking",   "STRATEGY"),
            ("content",     "Report Content Summary",       "CONTENT"),
            ("editor",      "Quality Assessment",           "QA"),
        ],
        "proposal": [
            ("research",    "Country Evidence Pack",        "RESEARCH"),
            ("analyst",     "Needs & Gap Analysis",         "ANALYSIS"),
            ("strategist",  "Strategic Priorities",         "STRATEGY"),
            ("proj_design", "Project Concepts",             "DESIGN"),
            ("costing",     "Budget & Results Framework",   "COSTING"),
            ("content",     "Proposal Content Summary",     "CONTENT"),
            ("editor",      "Quality Assessment",           "QA"),
        ],
        "external": [
            ("ext_assessor","Project Scoring & Assessment", "ASSESSMENT"),
            ("editor",      "Quality Assessment",           "QA"),
        ],
    }

    for slide_num, (agent_id, title, tag) in enumerate(
            slide_defs.get(workflow, slide_defs["needs"]), start=2):
        text    = agent_outputs.get(agent_id, f"Output from {agent_id} agent.")
        bullets = extract_bullets(text)
        add_content_slide(prs, slide_num, title, tag, bullets)

    # ── Final slide: Next Steps ───────────────────────────────────────────────
    last_num = len(slide_defs.get(workflow, slide_defs["needs"])) + 2
    next_steps = [
        "Review full agent outputs in the agent_outputs/ folder for detailed findings.",
        f"Validate key data points with in-country contacts and latest national sources.",
        "Share with relevant team members for inputs before stakeholder engagement.",
        "Update the CAD pipeline dashboard with new project entries as appropriate.",
        "Schedule follow-up review meeting to agree on priority interventions.",
    ]
    add_content_slide(prs, last_num, "Next Steps", "ACTION", next_steps)

    prs.save(str(out_path))
    print(f"[OK] PPT saved: {fname} ({out_path.stat().st_size // 1024}KB)")
    return out_path

# ══════════════════════════════════════════════════════════════════════════════
# EMAIL BUILDER
# ══════════════════════════════════════════════════════════════════════════════

def build_result_html(ctx, agent_outputs, verdict):
    country  = ctx["country"]
    sector   = ctx["sector"]
    workflow = ctx["workflow"]
    agents   = WORKFLOW_AGENTS.get(workflow, NEEDS_AGENTS)
    now      = datetime.now().strftime("%A, %d %B %Y %H:%M")

    vc  = {"READY":"#2DB88A", "READY WITH CAVEATS":"#AD833B", "NOT READY":"#9B120B"}
    vbg = {"READY":"#E6F7F2", "READY WITH CAVEATS":"#FFF3CD", "NOT READY":"#FDECEA"}
    color = vc.get(verdict, "#2DB88A")
    bg    = vbg.get(verdict, "#E6F7F2")

    agent_rows = "".join(
        f'<tr>'
        f'<td style="padding:7px 12px;border-bottom:1px solid #E5E8EF;font-size:12px;'
        f'color:#111827;font-weight:600">{aid}</td>'
        f'<td style="padding:7px 12px;border-bottom:1px solid #E5E8EF;font-size:11px;'
        f'color:#6B7280">{agent_outputs.get(aid,"")[:120].replace(chr(10)," ")}...</td>'
        f'</tr>'
        for aid in agents if aid in agent_outputs
    )

    return f"""<!DOCTYPE html>
<html><head><meta charset="UTF-8"></head>
<body style="margin:0;padding:0;background:#F4F5F8;font-family:Arial,sans-serif">
<table width="100%" cellpadding="0" cellspacing="0" style="background:#F4F5F8;padding:32px 0">
<tr><td align="center">
<table width="620" cellpadding="0" cellspacing="0"
  style="background:#fff;border-radius:12px;overflow:hidden;box-shadow:0 4px 24px rgba(0,0,0,.08)">

  <tr><td style="background:#2D3F7B;padding:24px 28px">
    <div style="font-size:18px;font-weight:800;color:#fff">Country Assessment Dashboard</div>
    <div style="font-size:13px;color:rgba(255,255,255,.7);margin-top:4px">
      Agent Pipeline Result &mdash; {now}
    </div>
  </td></tr>

  <tr><td style="padding:28px">

    <div style="background:{bg};border-radius:10px;padding:16px 20px;margin-bottom:24px;
                border-left:4px solid {color}">
      <div style="font-size:15px;font-weight:800;color:{color}">{verdict}</div>
      <div style="font-size:12px;color:#374151;margin-top:4px">
        <strong>{country}</strong> &middot; {sector} &middot; {workflow.upper()} workflow
        &middot; {len(agents)} agents
      </div>
    </div>

    <p style="font-size:13px;color:#374151;margin:0 0 20px">
      Full agent outputs are saved to <strong>agent_outputs/</strong> in your project folder.
    </p>

    <div style="margin-top:8px">
      <div style="font-size:13px;font-weight:700;color:#111827;margin-bottom:10px">
        Agent outputs summary
      </div>
      <table width="100%" cellpadding="0" cellspacing="0"
        style="border-collapse:collapse;border:1px solid #E5E8EF;border-radius:8px;
               overflow:hidden;font-family:Arial,sans-serif">
        <thead>
          <tr style="background:#F4F5F8">
            <th style="padding:8px 12px;text-align:left;font-size:11px;color:#6B7280;
                       font-weight:600;width:140px">AGENT</th>
            <th style="padding:8px 12px;text-align:left;font-size:11px;color:#6B7280;
                       font-weight:600">EXCERPT</th>
          </tr>
        </thead>
        <tbody>{agent_rows}</tbody>
      </table>
    </div>

    <div style="margin-top:20px;padding:14px 16px;background:#F4F5F8;border-radius:8px">
      <div style="font-size:11px;color:#6B7280;line-height:1.6">
        To trigger another run, send an email with subject:<br>
        <code>RUN: &lt;country&gt; / &lt;sector&gt; / &lt;workflow&gt;</code>
      </div>
    </div>

  </td></tr>

  <tr><td style="background:#F4F5F8;padding:16px 28px;border-top:1px solid #E5E8EF">
    <p style="margin:0;font-size:11px;color:#9CA3AF">
      Generated by the CAD Agent Pipeline &middot; Office of Development Affairs<br>
      {len(agents)} specialist agents &middot; claude-opus-4-5
    </p>
  </td></tr>

</table></td></tr></table>
</body></html>"""


# ══════════════════════════════════════════════════════════════════════════════
# EMAIL SEND
# ══════════════════════════════════════════════════════════════════════════════

def send_result_email(ctx, agent_outputs, verdict, reply_to_address, reply_subject, pptx_path=None):
    sender   = ENV.get("GMAIL_SENDER", "")
    password = ENV.get("GMAIL_APP_PASSWORD", "").replace(" ", "")

    if not sender or not password:
        print("[WARN] No Gmail credentials -- saving preview only.")
        return False

    html_body = build_result_html(ctx, agent_outputs, verdict)
    subject   = f"RE: {reply_subject} -- {verdict}"

    # Mixed message so we can attach the PPT
    msg = MIMEMultipart("mixed")
    msg["Subject"] = subject
    msg["From"]    = f"CAD Agent Pipeline <{sender}>"
    msg["To"]      = reply_to_address

    # HTML body
    alt = MIMEMultipart("alternative")
    alt.attach(MIMEText(html_body, "html"))
    msg.attach(alt)

    # Attach PPT if available
    if pptx_path and Path(pptx_path).exists():
        with open(pptx_path, "rb") as f:
            part = MIMEBase("application",
                            "vnd.openxmlformats-officedocument.presentationml.presentation")
            part.set_payload(f.read())
        encoders.encode_base64(part)
        part.add_header("Content-Disposition", "attachment",
                        filename=Path(pptx_path).name)
        msg.attach(part)
        size_kb = Path(pptx_path).stat().st_size // 1024
        print(f"[OK] Attaching {Path(pptx_path).name} ({size_kb}KB)")
    else:
        print("[WARN] No PPT to attach")

    try:
        with smtplib.SMTP_SSL("smtp.gmail.com", 465) as smtp:
            smtp.login(sender, password)
            smtp.sendmail(sender, reply_to_address, msg.as_string())
        print(f"[OK] Result email sent to {reply_to_address}")
        return True
    except smtplib.SMTPAuthenticationError:
        print("[ERR] Gmail auth failed -- check App Password in .env")
        preview = OUT_DIR / f"{ctx['country'].replace(' ','_')}_{ctx['workflow']}_result.html"
        preview.write_text(html_body, encoding="utf-8")
        print(f"[INFO] Preview saved: {preview}")
    except Exception as e:
        print(f"[ERR] Email send failed: {e}")
    return False


# ══════════════════════════════════════════════════════════════════════════════
# INBOX POLLER
# ══════════════════════════════════════════════════════════════════════════════

def parse_trigger_subject(subject):
    if not subject.upper().startswith(TRIGGER_PREFIX.upper()):
        return None
    body  = subject[len(TRIGGER_PREFIX):].strip()
    parts = [p.strip() for p in body.split("/")]
    if len(parts) < 3:
        return None
    workflow = parts[2].lower().strip()
    if workflow not in WORKFLOW_AGENTS:
        if "needs" in workflow or "report" in workflow:       workflow = "needs"
        elif "proposal" in workflow or "program" in workflow: workflow = "proposal"
        elif "external" in workflow or "assess" in workflow:  workflow = "external"
        else:                                                  workflow = "needs"
    return {
        "country":  parts[0],
        "sector":   parts[1],
        "workflow": workflow,
        "audience": "Senior Leadership & Donor Partners",
        "budget":   "TBD",
        "horizon":  "3 years",
    }


def _parse_email_addr(raw):
    """Pull bare email out of 'Name <addr@host>' or 'addr@host'."""
    if not raw:
        return ""
    m = re.search(r"<(.+?)>", raw)
    return (m.group(1) if m else raw).strip().lower()


def _allowed_senders():
    """Comma-separated list in env ALLOWED_SENDERS (lowercased)."""
    raw = ENV.get("ALLOWED_SENDERS", "").strip()
    if not raw:
        return set()
    return {a.strip().lower() for a in raw.split(",") if a.strip()}


def check_inbox_for_triggers(seen_ids):
    import imaplib, email as email_lib

    sender   = ENV.get("GMAIL_SENDER", "")
    password = ENV.get("GMAIL_APP_PASSWORD", "").replace(" ", "")
    if not sender or not password:
        return []

    # ── Sender whitelist (deny-by-default) ──
    # An empty ALLOWED_SENDERS means: only the inbox owner may trigger the
    # pipeline. Anyone else's RUN: email is marked Seen and discarded so the
    # Anthropic spend is gated to authorised actors.
    allowed = _allowed_senders()
    if not allowed:
        allowed = {sender.lower()}
        print(f"[info] ALLOWED_SENDERS empty — defaulting to inbox owner only ({sender})")

    triggers = []
    try:
        with imaplib.IMAP4_SSL("imap.gmail.com") as mail:
            mail.login(sender, password)
            mail.select("INBOX")
            _, msg_nums = mail.search(None, f'(UNSEEN SUBJECT "{TRIGGER_PREFIX}")')

            for num in msg_nums[0].split():
                _, data = mail.fetch(num, "(RFC822)")
                raw = data[0][1]
                msg = email_lib.message_from_bytes(raw)

                msg_id   = msg.get("Message-ID", "")
                subject  = msg.get("Subject", "")
                from_hdr = msg.get("From", "")
                reply_to = msg.get("Reply-To") or from_hdr

                from_addr  = _parse_email_addr(from_hdr)
                reply_addr = _parse_email_addr(reply_to)

                if msg_id in seen_ids:
                    continue

                # Sender authorization gate
                if from_addr not in allowed:
                    print(f"[deny] Trigger from unauthorized sender {from_addr!r} "
                          f"— subject={subject!r}")
                    # Mark Seen so we don't re-process repeatedly
                    mail.store(num, "+FLAGS", "\\Seen")
                    seen_ids.add(msg_id)
                    continue

                # Reply-To spoof guard: if Reply-To differs from authenticated From,
                # send results back to From (the authorised sender), not the spoofed
                # Reply-To. Prevents exfiltration of the generated PPT.
                if reply_addr != from_addr:
                    print(f"[warn] Reply-To {reply_addr!r} != From {from_addr!r} — "
                          f"replying to From only")
                    reply_addr = from_addr

                ctx = parse_trigger_subject(subject)
                if ctx:
                    mail.store(num, "+FLAGS", "\\Seen")
                    triggers.append((ctx, reply_addr, subject, msg_id))
                    seen_ids.add(msg_id)
                    print(f"[TRIGGER] {subject} -> from {reply_addr}")

    except imaplib.IMAP4.error as e:
        print(f"[WARN] IMAP login failed: {e}")
        print("       Check GMAIL_APP_PASSWORD in .env (must be Google App Password, not regular password)")
    except Exception as e:
        print(f"[WARN] Inbox check failed: {e}")

    return triggers


def run_once():
    """Single inbox check — used by GitHub Actions (run and exit)."""
    print(f"[{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}] CAD Pipeline — single check")
    seen_ids = set()
    triggers = check_inbox_for_triggers(seen_ids)
    if not triggers:
        print("No triggers found.")
        return
    print(f"{len(triggers)} trigger(s) found!")
    for ctx, reply_to, reply_subject, msg_id in triggers:
        print(f"\n-> Running: {ctx['country']} / {ctx['sector']} / {ctx['workflow']}")
        try:
            outputs = run_pipeline(ctx)
            editor  = outputs.get("editor", "")
            verdict = ("NOT READY"          if "NOT READY" in editor.upper()
                       else "READY WITH CAVEATS" if "CAVEATS" in editor.upper()
                       else "READY")
            print(f"\nVerdict: {verdict}")
            print("Generating PowerPoint...")
            pptx_path = generate_pptx(ctx, outputs, verdict)
            html    = build_result_html(ctx, outputs, verdict)
            preview = OUT_DIR / f"{ctx['country'].replace(' ','_')}_{ctx['workflow']}_result.html"
            preview.write_text(html, encoding="utf-8")
            send_result_email(ctx, outputs, verdict, reply_to, reply_subject, pptx_path=pptx_path)
        except Exception as e:
            print(f"[ERR] Pipeline failed: {e}")


def run_watch_mode():
    print(f"\n{'='*56}")
    print(f"  CAD Agent Pipeline -- WATCH MODE")
    print(f"  Polling inbox every {POLL_INTERVAL_SECONDS//60} min")
    print(f"  Trigger: send email with subject starting '{TRIGGER_PREFIX}'")
    print(f"  Format:  RUN: <country> / <sector> / <workflow>")
    print(f"  To stop: Ctrl+C or close this window")
    print(f"{'='*56}\n")

    seen_ids = set()
    while True:
        print(f"[{datetime.now().strftime('%H:%M:%S')}] Checking inbox...", end=" ", flush=True)
        triggers = check_inbox_for_triggers(seen_ids)

        if not triggers:
            print("no triggers found.")
        else:
            print(f"{len(triggers)} trigger(s) found!")
            for ctx, reply_to, reply_subject, msg_id in triggers:
                print(f"\n  -> Running: {ctx['country']} / {ctx['sector']} / {ctx['workflow']}")
                try:
                    outputs = run_pipeline(ctx)
                    editor  = outputs.get("editor", "")
                    verdict = ("NOT READY"          if "NOT READY" in editor.upper()
                               else "READY WITH CAVEATS" if "CAVEATS" in editor.upper()
                               else "READY")
                    print(f"\n  Verdict: {verdict}")
                    print("  Generating PowerPoint...")
                    pptx_path = generate_pptx(ctx, outputs, verdict)
                    html    = build_result_html(ctx, outputs, verdict)
                    preview = OUT_DIR / f"{ctx['country'].replace(' ','_')}_{ctx['workflow']}_result.html"
                    preview.write_text(html, encoding="utf-8")
                    print(f"  Preview saved: {preview.name}")
                    send_result_email(ctx, outputs, verdict, reply_to, reply_subject,
                                      pptx_path=pptx_path)
                except Exception as e:
                    print(f"[ERR] Pipeline failed: {e}")

        print(f"  Next check in {POLL_INTERVAL_SECONDS//60} min...")
        time.sleep(POLL_INTERVAL_SECONDS)


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════

def main():
    args = sys.argv[1:]

    if "--watch" in args:
        run_watch_mode()
        return
    if "--once" in args:
        run_once()
        return

    # Manual run defaults -- edit these or pass via CLI flags
    ctx = {
        "country":  "Jordan",
        "sector":   "Education & Youth Employment",
        "workflow": "needs",
        "audience": "Senior Leadership & Donor Partners",
        "budget":   "TBD",
        "horizon":  "3 years",
    }

    for flag, key in [("--country","country"), ("--sector","sector"),
                      ("--workflow","workflow"), ("--budget","budget")]:
        if flag in args:
            idx = args.index(flag)
            if idx + 1 < len(args):
                ctx[key] = args[idx + 1]

    outputs = run_pipeline(ctx)

    editor  = outputs.get("editor", "")
    verdict = ("NOT READY"          if "NOT READY" in editor.upper()
               else "READY WITH CAVEATS" if "CAVEATS" in editor.upper()
               else "READY")

    print(f"\n  Final verdict: {verdict}")

    # Generate PPT
    print("\n  Generating PowerPoint...")
    pptx_path = generate_pptx(ctx, outputs, verdict)

    # Save HTML preview
    html    = build_result_html(ctx, outputs, verdict)
    preview = OUT_DIR / f"{ctx['country'].replace(' ','_')}_{ctx['workflow']}_result.html"
    preview.write_text(html, encoding="utf-8")
    print(f"[OK] Preview saved: {preview}")

    if "--email" in args:
        recipient = ENV.get("DIGEST_RECIPIENT") or ENV.get("GMAIL_SENDER", "")
        if recipient:
            send_result_email(ctx, outputs, verdict, recipient,
                              f"RUN: {ctx['country']} / {ctx['sector']} / {ctx['workflow']}",
                              pptx_path=pptx_path)


if __name__ == "__main__":
    main()
