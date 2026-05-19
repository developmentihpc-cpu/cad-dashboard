"""
country_ppt_builder.py
======================
SINGLE SOURCE OF TRUTH for the CAD 13-slide country brief.

Design reference:
    design/Country_Overview_extracted.html
    (canonical layout — keep this builder in sync with that file)

Called by:
  - generate_ethiopia_ppt.py   (one-off chat scripts)
  - agents_pipeline.py         (email-triggered agents)
  - telegram_webhook.py        (Cloud Run /generate-ppt endpoint)

Edit this file and ALL three paths inherit the change.

Entry point:
    build(context, output=None) -> bytes (if output is None) or None

Context schema: see CONTEXT_SCHEMA at bottom for full reference.
"""

from __future__ import annotations
import io
import os
import urllib.request
from datetime import datetime


def _coerce_color(v):
    """Accept RGBColor, [r,g,b] list/tuple, or '#RRGGBB' hex; return RGBColor or pass through."""
    from pptx.dml.color import RGBColor as _RGB
    if isinstance(v, _RGB):
        return v
    if isinstance(v, str) and v.startswith("#") and len(v) == 7:
        try:
            return _RGB(int(v[1:3], 16), int(v[3:5], 16), int(v[5:7], 16))
        except Exception:
            return v
    if isinstance(v, (list, tuple)) and len(v) == 3 and all(isinstance(x, int) for x in v):
        try:
            return _RGB(*v)
        except Exception:
            return v
    return v


# Context keys that hold colors — coerced at build() entry.
_COLOR_KEYS = {"color", "value_color", "bg", "bar_color", "tag_color",
               "fg", "stroke", "accent"}


def _safe_copy(obj, *, _color_key=False):
    """deepcopy-equivalent that:
       - keeps RGBColor instances intact (RGBColor is a tuple subclass, so the
         generic tuple branch would degrade them back to plain tuples)
       - coerces values keyed by 'color', 'bg', etc. from JSON-friendly forms
         (list/tuple of ints, hex string) back to RGBColor."""
    from pptx.dml.color import RGBColor as _RGB
    if isinstance(obj, _RGB):
        return obj
    if _color_key:
        coerced = _coerce_color(obj)
        if isinstance(coerced, _RGB):
            return coerced
    if isinstance(obj, dict):
        return {k: _safe_copy(v, _color_key=(k in _COLOR_KEYS))
                for k, v in obj.items()}
    if isinstance(obj, list):
        return [_safe_copy(v, _color_key=_color_key) for v in obj]
    if isinstance(obj, tuple):
        return tuple(_safe_copy(v, _color_key=_color_key) for v in obj)
    return obj

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ─── Palette (matches design/Country_Overview_extracted.html :root) ───────────
GOLD       = RGBColor(0xAD, 0x83, 0x3B)
GOLD_LT    = RGBColor(0xC7, 0xA8, 0x77)
GOLD_PALE  = RGBColor(0xF5, 0xEF, 0xE3)
BLUE       = RGBColor(0x67, 0x8C, 0xA5)
BLUE_LT    = RGBColor(0xCB, 0xDC, 0xE6)
BLUE_PALE  = RGBColor(0xEE, 0xF4, 0xF8)
BLACK      = RGBColor(0x10, 0x18, 0x20)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DEEP_RED   = RGBColor(0x79, 0x24, 0x2F)
DEEP_GREEN = RGBColor(0x2F, 0x58, 0x6E)
NAVY       = RGBColor(0x33, 0x3F, 0x64)
NAVY_DK    = RGBColor(0x1A, 0x21, 0x38)  # cover hero gradient endpoint
POS        = RGBColor(0x00, 0xBC, 0x8B)
SEMI_POS   = RGBColor(0x9C, 0xBB, 0x5D)
NEUTRAL    = RGBColor(0xF3, 0x9B, 0x26)
SEMI_NEG   = RGBColor(0xBA, 0x49, 0x2F)
NEG        = RGBColor(0x9B, 0x12, 0x0B)
FG         = RGBColor(0x1C, 0x24, 0x33)
FG2        = RGBColor(0x5B, 0x6A, 0x7E)
FG3        = RGBColor(0x88, 0x96, 0xA6)
BG         = RGBColor(0xF4, 0xF1, 0xEA)  # warm cream
BORDER     = RGBColor(0xDC, 0xE3, 0xEA)
RULE       = RGBColor(0xC9, 0xA5, 0x5A)
MAP_BG     = RGBColor(0xD8, 0xE7, 0xF0)

# Status chip colors
STATUS_COLORS = {
    "improving":  POS,
    "developing": GOLD,
    "weak":       SEMI_NEG,
    "fragile":    SEMI_NEG,
    "severe":     NEG,
}

# python-pptx uses Calibri as fallback; PowerPoint substitutes DM Sans if installed
FONT  = "DM Sans"
FONT_FALLBACK = "Calibri"

# 16:9 widescreen at 13.333" × 7.5" (matches design's 1920×1080 ratio)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ─── Low-level primitives ─────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill, line=None, line_w_pt=0.5):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w_pt)
    s.shadow.inherit = False
    return s


def add_oval(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
    s.shadow.inherit = False
    return s


def add_line(slide, x, y, w, h, color, *, weight_pt=0.5, dashed=False):
    if dashed:
        ln = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
                                        x, y, x + w, y + h)
        ln.line.color.rgb = color
        ln.line.width = Pt(weight_pt)
        lnPr = ln.line._get_or_add_ln()
        prstDash = etree.SubElement(lnPr, qn('a:prstDash'))
        prstDash.set('val', 'dash')
        return ln
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def add_text(slide, x, y, w, h, text, *,
             size=11, bold=False, color=FG, font=FONT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             char_spacing=None, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = Inches(0.03)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    if char_spacing is not None:
        rPr = r._r.get_or_add_rPr()
        rPr.set('spc', str(int(char_spacing * 100)))
    return tb


def add_runs(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, wrap=True):
    """Multi-run textbox. runs = list of dicts: {text, size, bold, color, italic, font, spc}."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = Inches(0.03)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    for r_def in runs:
        r = p.add_run()
        r.text = r_def.get("text", "")
        r.font.name = r_def.get("font", FONT)
        r.font.size = Pt(r_def.get("size", 11))
        r.font.bold = r_def.get("bold", False)
        r.font.italic = r_def.get("italic", False)
        if "color" in r_def:
            r.font.color.rgb = r_def["color"]
        if "spc" in r_def:
            rPr = r._r.get_or_add_rPr()
            rPr.set('spc', str(int(r_def["spc"] * 100)))
    return tb


# ─── Flag fetch ────────────────────────────────────────────────────────────────
def _decode_image_bytes(v) -> bytes | None:
    """Accept raw bytes OR base64 string (for JSON transport) OR None."""
    if v is None:
        return None
    if isinstance(v, (bytes, bytearray)):
        return bytes(v)
    if isinstance(v, str):
        try:
            import base64
            # Strip any data: URI prefix
            s = v.split(",", 1)[1] if v.startswith("data:") else v
            return base64.b64decode(s)
        except Exception:
            return None
    return None


def _ctx_map_image(ctx: dict) -> bytes | None:
    """Pull the map image bytes from ctx['map']['image_bytes'], if present.
    Set once at build() entry by the renderer; reused by every slide that
    has a map placeholder."""
    m = ctx.get("map") or {}
    return _decode_image_bytes(m.get("image_bytes"))


_FLAG_CACHE: dict[str, bytes | None] = {}

def fetch_flag(iso2: str, custom_url: str | None = None) -> bytes | None:
    key = (custom_url or iso2 or "").lower()
    if key in _FLAG_CACHE:
        return _FLAG_CACHE[key]
    url = custom_url or (f"https://flagcdn.com/w160/{iso2.lower()}.png" if iso2 else None)
    if not url:
        _FLAG_CACHE[key] = None
        return None
    try:
        with urllib.request.urlopen(url, timeout=4) as r:
            data = r.read()
            _FLAG_CACHE[key] = data
            return data
    except Exception:
        _FLAG_CACHE[key] = None
        return None


# ─── Shared blocks ────────────────────────────────────────────────────────────

def add_header(slide, crumb, topic, statement, *, crumb_color=BLUE,
               topic_color=NAVY, stmt_color=GOLD):
    """Standard slide header — used on slides 2–10."""
    # Crumb
    add_text(slide, Inches(0.55), Inches(0.30), Inches(12.2), Inches(0.26),
             crumb.upper(), size=11, bold=True, color=crumb_color, char_spacing=2.5)
    # Two-tone title
    add_runs(slide, Inches(0.55), Inches(0.58), Inches(12.2), Inches(0.95),
             [
                 {"text": topic + " ",   "size": 30, "bold": True, "color": topic_color},
                 {"text": "|  ",         "size": 30, "color": GOLD},
                 {"text": statement,     "size": 30, "bold": True, "color": stmt_color},
             ], anchor=MSO_ANCHOR.MIDDLE)


def add_footer(slide, refs_text):
    """Numbered superscript sources footer."""
    add_line(slide, Inches(0.55), Inches(7.02), Inches(12.2), Emu(4000), BORDER)
    add_text(slide, Inches(0.55), Inches(7.10), Inches(12.2), Inches(0.30),
             refs_text, size=9, color=FG3, italic=True)


def add_section_label(slide, x, y, w, text):
    """Subheading above a viz / stats column."""
    add_text(slide, x, y, w, Inches(0.30),
             text, size=12, bold=True, color=FG2, char_spacing=1)
    add_line(slide, x, y + Inches(0.32), Inches(1.0), Emu(5000), GOLD, weight_pt=1.0)


def add_stat_row(slide, x, y, w, h, value, label, *, value_color=GOLD,
                 emphasize_neg=False):
    """Big-number + descriptive label row used on every sector slide."""
    if emphasize_neg:
        value_color = NEG
    add_text(slide, x, y, Inches(1.7), h,
             value, size=30, bold=True, color=value_color,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x + Inches(1.75), y, w - Inches(1.75), h,
             label, size=10.5, color=FG, anchor=MSO_ANCHOR.MIDDLE)


def add_stats_block(slide, x, y, w, items):
    """items: list of {val, label, neg?}"""
    rH = Inches(0.72)
    for i, it in enumerate(items[:4]):
        ry = y + rH * i
        if i > 0:
            add_line(slide, x, ry, w, Emu(3000), BORDER, weight_pt=0.4, dashed=True)
        add_stat_row(slide, x, ry + Inches(0.06), w, rH - Inches(0.10),
                     str(it.get("val", "—")),
                     str(it.get("label", "")),
                     emphasize_neg=it.get("neg", False))


def add_insights_block(slide, x, y, w, h, items, *, title="Key Insights",
                       bg=BLUE_PALE, bar_color=BLUE):
    add_rect(slide, x, y, Inches(0.06), h, bar_color)
    add_rect(slide, x + Inches(0.06), y, w - Inches(0.06), h, bg)
    add_text(slide, x + Inches(0.22), y + Inches(0.12),
             w - Inches(0.30), Inches(0.26),
             title.upper(), size=10, bold=True, color=bar_color, char_spacing=2)
    if not items:
        return
    items = items[:4]
    ih = (h - Inches(0.44)) / len(items)
    for i, item in enumerate(items):
        ty = y + Inches(0.42) + ih * i
        add_oval(slide, x + Inches(0.22), ty + Inches(0.10),
                 Inches(0.07), Inches(0.07), bar_color)
        add_text(slide, x + Inches(0.34), ty,
                 w - Inches(0.42), ih,
                 item, size=10, color=FG, anchor=MSO_ANCHOR.TOP)


def add_intv_block(slide, x, y, w, h, items, *, title="Potential Interventions"):
    add_rect(slide, x, y, Inches(0.06), h, GOLD)
    add_rect(slide, x + Inches(0.06), y, w - Inches(0.06), h, GOLD_PALE)
    add_text(slide, x + Inches(0.22), y + Inches(0.10),
             w - Inches(0.30), Inches(0.24),
             title.upper(), size=9.5, bold=True, color=GOLD, char_spacing=2)
    if not items:
        return
    items = items[:4]
    ih = (h - Inches(0.40)) / len(items)
    for i, item in enumerate(items):
        ty = y + Inches(0.38) + ih * i
        add_oval(slide, x + Inches(0.22), ty + Inches(0.10),
                 Inches(0.07), Inches(0.07), GOLD)
        add_text(slide, x + Inches(0.34), ty,
                 w - Inches(0.42), ih,
                 item, size=10, color=FG, anchor=MSO_ANCHOR.TOP)


# ─── Visualization helpers (slide-specific) ────────────────────────────────────

def add_viz_frame(slide, x, y, w, h, title=None):
    add_rect(slide, x, y, w, h, WHITE, line=BORDER)
    if title:
        add_text(slide, x + Inches(0.20), y + Inches(0.16),
                 w - Inches(0.40), Inches(0.24),
                 title, size=11, bold=True, color=FG2, char_spacing=1)
    return (x + Inches(0.20), y + Inches(0.50),
            w - Inches(0.40), h - Inches(0.70))


def add_bar_line_chart(slide, x, y, w, h, *, title, x_labels, bar_vals,
                       line_vals, y_max=40, y_step=10, bar_label="GDP growth",
                       line_label="Inflation"):
    """Bars + overlay line (used on slide 3 Economy)."""
    inner_x, inner_y, inner_w, inner_h = add_viz_frame(slide, x, y, w, h, title)
    chart_x = inner_x + Inches(0.4)
    chart_y = inner_y + Inches(0.10)
    chart_w = inner_w - Inches(0.4)
    chart_h = inner_h - Inches(0.60)
    # y-axis grid lines
    n_grid = int(y_max / y_step)
    for i in range(n_grid + 1):
        yy = chart_y + chart_h * i / n_grid
        add_line(slide, chart_x, yy, chart_w, Emu(2000),
                 BORDER if i > 0 else FG3,
                 weight_pt=0.5, dashed=(i > 0))
        add_text(slide, chart_x - Inches(0.4), yy - Inches(0.10),
                 Inches(0.35), Inches(0.22),
                 f"{int(y_max - i * y_step)}%", size=9, color=FG2,
                 align=PP_ALIGN.RIGHT)
    # Bars
    n = len(x_labels)
    col_w = chart_w / n
    bar_w = col_w * 0.50
    for i, v in enumerate(bar_vals):
        bh = chart_h * (v / y_max)
        bx = chart_x + col_w * i + (col_w - bar_w) / 2
        by = chart_y + chart_h - bh
        col = GOLD if i == n - 1 else BLUE
        add_rect(slide, bx, by, bar_w, bh, col)
        add_text(slide, bx - Inches(0.2), by - Inches(0.28),
                 bar_w + Inches(0.4), Inches(0.24),
                 f"{v}%", size=10, bold=True,
                 color=GOLD if i == n - 1 else FG, align=PP_ALIGN.CENTER)
        add_text(slide, chart_x + col_w * i, chart_y + chart_h + Inches(0.05),
                 col_w, Inches(0.24),
                 x_labels[i], size=9, color=FG2, align=PP_ALIGN.CENTER)
    # Line overlay (inflation)
    pts = []
    for i, v in enumerate(line_vals):
        cx = chart_x + col_w * i + col_w / 2
        cy = chart_y + chart_h * (1 - v / y_max)
        pts.append((cx, cy))
    for i in range(len(pts) - 1):
        # short rotated rectangle as connector
        from math import atan2, degrees, hypot
        x1, y1 = pts[i]; x2, y2 = pts[i + 1]
        # use straight connector for simplicity
        ln = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
                                        x1, y1, x2, y2)
        ln.line.color.rgb = NEG
        ln.line.width = Pt(2.5)
    for cx, cy in pts:
        add_oval(slide, cx - Inches(0.08), cy - Inches(0.08),
                 Inches(0.16), Inches(0.16), NEG)
    if line_vals:
        last_x, last_y = pts[-1]
        add_text(slide, last_x - Inches(0.40), last_y - Inches(0.35),
                 Inches(0.80), Inches(0.24),
                 f"{line_vals[-1]}%", size=11, bold=True,
                 color=NEG, align=PP_ALIGN.CENTER)
    # Legend
    leg_y = inner_y + inner_h - Inches(0.10)
    add_rect(slide, chart_x + Inches(0.5), leg_y,
             Inches(0.18), Inches(0.18), GOLD)
    add_text(slide, chart_x + Inches(0.72), leg_y - Inches(0.04),
             Inches(1.6), Inches(0.26),
             bar_label, size=9.5, color=FG)
    add_oval(slide, chart_x + Inches(2.7), leg_y + Inches(0.04),
             Inches(0.12), Inches(0.12), NEG)
    add_text(slide, chart_x + Inches(2.9), leg_y - Inches(0.04),
             Inches(1.8), Inches(0.26),
             line_label, size=9.5, color=FG)


def add_grouped_bars(slide, x, y, w, h, *, title, x_labels, series_a, series_b,
                     a_label, b_label, a_color=BLUE, b_color=GOLD,
                     y_max=120, y_step=20):
    """Two grouped bar series (used on slide 5 Education)."""
    inner_x, inner_y, inner_w, inner_h = add_viz_frame(slide, x, y, w, h, title)
    chart_x = inner_x + Inches(0.4)
    chart_y = inner_y + Inches(0.10)
    chart_w = inner_w - Inches(0.4)
    chart_h = inner_h - Inches(0.60)
    n_grid = int(y_max / y_step)
    for i in range(n_grid + 1):
        yy = chart_y + chart_h * i / n_grid
        add_line(slide, chart_x, yy, chart_w, Emu(2000),
                 BORDER if i > 0 else FG3,
                 weight_pt=0.5, dashed=(i > 0))
        add_text(slide, chart_x - Inches(0.45), yy - Inches(0.10),
                 Inches(0.40), Inches(0.22),
                 f"{int(y_max - i * y_step)}%", size=9, color=FG2,
                 align=PP_ALIGN.RIGHT)
    n = len(x_labels)
    grp_w = chart_w / n
    bar_w = grp_w * 0.32
    gap = grp_w * 0.06
    for i, (a, b) in enumerate(zip(series_a, series_b)):
        ah = chart_h * (a / y_max); bh = chart_h * (b / y_max)
        ax = chart_x + grp_w * i + (grp_w / 2) - bar_w - gap / 2
        bx = chart_x + grp_w * i + (grp_w / 2) + gap / 2
        ay = chart_y + chart_h - ah; by = chart_y + chart_h - bh
        add_rect(slide, ax, ay, bar_w, ah, a_color)
        add_rect(slide, bx, by, bar_w, bh, b_color)
        add_text(slide, chart_x + grp_w * i, chart_y + chart_h + Inches(0.05),
                 grp_w, Inches(0.24),
                 x_labels[i], size=9, color=FG2, align=PP_ALIGN.CENTER)
    # Legend
    leg_y = inner_y + inner_h - Inches(0.10)
    add_rect(slide, chart_x + Inches(1.0), leg_y, Inches(0.18), Inches(0.18), a_color)
    add_text(slide, chart_x + Inches(1.22), leg_y - Inches(0.04),
             Inches(1.5), Inches(0.26), a_label, size=9.5, color=FG)
    add_rect(slide, chart_x + Inches(3.0), leg_y, Inches(0.18), Inches(0.18), b_color)
    add_text(slide, chart_x + Inches(3.22), leg_y - Inches(0.04),
             Inches(1.5), Inches(0.26), b_label, size=9.5, color=FG)


def add_line_chart(slide, x, y, w, h, *, title, x_labels, series,
                   y_max=50, y_step=10):
    """series: list of {label, values, color, dashed?}"""
    inner_x, inner_y, inner_w, inner_h = add_viz_frame(slide, x, y, w, h, title)
    chart_x = inner_x + Inches(0.4)
    chart_y = inner_y + Inches(0.10)
    chart_w = inner_w - Inches(0.4)
    chart_h = inner_h - Inches(0.60)
    n_grid = int(y_max / y_step)
    for i in range(n_grid + 1):
        yy = chart_y + chart_h * i / n_grid
        add_line(slide, chart_x, yy, chart_w, Emu(2000),
                 BORDER if i > 0 else FG3,
                 weight_pt=0.5, dashed=(i > 0))
        add_text(slide, chart_x - Inches(0.45), yy - Inches(0.10),
                 Inches(0.40), Inches(0.22),
                 str(int(y_max - i * y_step)), size=9, color=FG2,
                 align=PP_ALIGN.RIGHT)
    n = len(x_labels)
    col_w = chart_w / max(n - 1, 1)
    for s in series:
        pts = []
        for i, v in enumerate(s["values"]):
            cx = chart_x + col_w * i
            cy = chart_y + chart_h * (1 - v / y_max)
            pts.append((cx, cy))
        for i in range(len(pts) - 1):
            x1, y1 = pts[i]; x2, y2 = pts[i + 1]
            ln = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
                                            x1, y1, x2, y2)
            ln.line.color.rgb = s["color"]
            ln.line.width = Pt(2.5)
            if s.get("dashed"):
                lnPr = ln.line._get_or_add_ln()
                etree.SubElement(lnPr, qn('a:prstDash')).set('val', 'dash')
        for cx, cy in pts:
            add_oval(slide, cx - Inches(0.07), cy - Inches(0.07),
                     Inches(0.14), Inches(0.14), s["color"])
    # x labels
    for i, lbl in enumerate(x_labels):
        cx = chart_x + col_w * i - Inches(0.30)
        add_text(slide, cx, chart_y + chart_h + Inches(0.05),
                 Inches(0.60), Inches(0.24),
                 lbl, size=9, color=FG2, align=PP_ALIGN.CENTER)
    # legend
    leg_y = inner_y + inner_h - Inches(0.10)
    cur_x = chart_x + Inches(0.5)
    for s in series:
        add_oval(slide, cur_x, leg_y + Inches(0.03),
                 Inches(0.13), Inches(0.13), s["color"])
        add_text(slide, cur_x + Inches(0.18), leg_y - Inches(0.04),
                 Inches(2.6), Inches(0.26),
                 s["label"], size=9.5, color=FG)
        cur_x += Inches(2.8)


def add_funnel(slide, x, y, w, h, *, title, rows):
    """Population cascade — used on slide 6 (Nutrition).
    rows: list of {pct, percent_text, label, sublabel, width_ratio, color}"""
    inner_x, inner_y, inner_w, inner_h = add_viz_frame(slide, x, y, w, h, title)
    if not rows:
        return
    rH = inner_h / len(rows)
    pct_col_w = Inches(1.4)
    stage_col_w = Inches(2.6)
    bar_col_w = inner_w - pct_col_w - stage_col_w - Inches(0.3)
    for i, r in enumerate(rows):
        ry = inner_y + rH * i + Inches(0.06)
        rh = rH - Inches(0.12)
        # Left: big pct + small percent_text
        add_text(slide, inner_x, ry, pct_col_w, rh,
                 r["pct"], size=22, bold=True, color=FG, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, inner_x, ry + rh - Inches(0.22),
                 pct_col_w, Inches(0.20),
                 r.get("percent_text", ""), size=9.5, color=FG3,
                 anchor=MSO_ANCHOR.MIDDLE)
        # Middle: bar
        bx = inner_x + pct_col_w + Inches(0.05)
        bw_actual = bar_col_w * r.get("width_ratio", 1.0)
        add_rect(slide, bx, ry + Inches(0.06), bw_actual, rh - Inches(0.12),
                 r["color"])
        add_text(slide, bx + Inches(0.15), ry + Inches(0.06),
                 bw_actual - Inches(0.20), rh - Inches(0.12),
                 r["label"], size=11, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)
        # Right: stage
        sx = inner_x + pct_col_w + bar_col_w + Inches(0.20)
        add_text(slide, sx, ry, stage_col_w, rh / 2,
                 r.get("stage", ""), size=10, bold=True, color=FG,
                 anchor=MSO_ANCHOR.BOTTOM)
        add_text(slide, sx, ry + rh / 2, stage_col_w, rh / 2,
                 r.get("sublabel", ""), size=8.5, color=FG3,
                 italic=True, anchor=MSO_ANCHOR.TOP)


def add_compare_rows(slide, x, y, w, h, *, title, rows):
    """Urban vs Rural comparison — slide 8.
    rows: [{label, urban, rural}]"""
    inner_x, inner_y, inner_w, inner_h = add_viz_frame(slide, x, y, w, h, title)
    # Column headers
    col_label_w  = inner_w * 0.30
    col_bar_w    = inner_w * 0.32
    col_value_w  = (inner_w * 0.38) / 2
    hx = inner_x + col_label_w
    add_text(slide, hx, inner_y, col_bar_w, Inches(0.30),
             "URBAN", size=11, bold=True, color=BLUE,
             align=PP_ALIGN.CENTER, char_spacing=2)
    add_text(slide, hx + col_bar_w + Inches(0.1), inner_y,
             col_bar_w - Inches(0.1), Inches(0.30),
             "RURAL", size=11, bold=True, color=NEG,
             align=PP_ALIGN.CENTER, char_spacing=2)
    rH = (inner_h - Inches(0.40)) / max(len(rows), 1)
    for i, row in enumerate(rows):
        ry = inner_y + Inches(0.40) + rH * i
        add_text(slide, inner_x, ry, col_label_w, rH,
                 row["label"], size=11, bold=True, color=FG,
                 anchor=MSO_ANCHOR.MIDDLE)
        # Urban bar (anchored right edge)
        u_pct = row["urban"] / 100
        u_w = col_bar_w * u_pct
        u_x = inner_x + col_label_w + col_bar_w - u_w
        add_rect(slide, u_x, ry + rH / 2 - Inches(0.16),
                 u_w, Inches(0.32), BLUE)
        add_text(slide, inner_x + col_label_w + col_bar_w - Inches(0.55),
                 ry + rH / 2 - Inches(0.16),
                 Inches(0.50), Inches(0.32),
                 f"{row['urban']}%", size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        # Rural bar (anchored left edge)
        r_pct = row["rural"] / 100
        r_w = col_bar_w * r_pct
        r_x = inner_x + col_label_w + col_bar_w + Inches(0.1)
        add_rect(slide, r_x, ry + rH / 2 - Inches(0.16),
                 r_w, Inches(0.32), NEG)
        add_text(slide, r_x + Inches(0.05), ry + rH / 2 - Inches(0.16),
                 Inches(0.50), Inches(0.32),
                 f"{row['rural']}%", size=10, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)


def add_map_placeholder(slide, x, y, w, h, *, country="", lat=None, lng=None,
                        legend=None, image_bytes=None):
    """Used on slides 1, 2, 4, 7, 10 — country/region map placeholder."""
    if image_bytes:
        try:
            slide.shapes.add_picture(io.BytesIO(image_bytes), x, y, w, h)
            return
        except Exception:
            pass
    add_rect(slide, x, y, w, h, MAP_BG, line=BORDER)
    # Simplified Ethiopia-like shape
    add_oval(slide, x + w * 0.30, y + h * 0.30, w * 0.45, h * 0.45, GOLD)
    add_text(slide, x, y + h * 0.45, w, Inches(0.40),
             (country or "MAP").upper(), size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, char_spacing=3)
    if legend:
        # Legend below
        lx = x + Inches(0.2); ly = y + h - Inches(0.50)
        for i, (lbl, col) in enumerate(legend):
            add_rect(slide, lx, ly, Inches(0.20), Inches(0.18), col)
            add_text(slide, lx + Inches(0.25), ly - Inches(0.04),
                     Inches(1.6), Inches(0.26),
                     lbl, size=9, color=FG)
            lx += Inches(2.0)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════════

def slide_1_cover(prs, ctx):
    """Cover: navy hero rail (left) + cream right with 3-col snapshot + map+sector status."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # Background — left navy rail, right cream
    LW = Inches(3.50)
    add_rect(s, 0, 0, LW, SLIDE_H, NAVY)
    add_rect(s, LW, 0, SLIDE_W - LW, SLIDE_H, BG)

    # Flag — either real image or 3-stripe placeholder
    flag = fetch_flag(ctx.get("iso2", ""), ctx.get("flag_url"))
    if flag:
        try:
            s.shapes.add_picture(io.BytesIO(flag),
                                 Inches(0.40), Inches(0.50),
                                 Inches(1.30), Inches(0.85))
        except Exception:
            flag = None
    if not flag:
        add_rect(s, Inches(0.40), Inches(0.50), Inches(1.30), Inches(0.85),
                 NAVY_DK, line=WHITE)

    add_text(s, Inches(0.40), Inches(1.55), LW - Inches(0.40), Inches(0.26),
             "COUNTRY OVERVIEW", size=11, bold=True, color=GOLD_LT, char_spacing=3)
    add_text(s, Inches(0.40), Inches(1.85), LW - Inches(0.50), Inches(1.20),
             ctx["country"], size=44, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.TOP)
    add_text(s, Inches(0.40), Inches(3.10), LW - Inches(0.50), Inches(0.60),
             ctx.get("subtitle", ""), size=11, color=GOLD_LT, italic=True)

    # Bottom metadata rows
    rows = [
        ("DAC CLASS",  ctx.get("snapshot", {}).get("dac_class", "—"), True),
        ("CAPITAL",    ctx.get("snapshot", {}).get("capital", "—"),    False),
        ("GNI / CAP",  ctx.get("snapshot", {}).get("gni_cap", "—"),    False),
        ("CURRENCY",   ctx.get("snapshot", {}).get("currency", "—"),   False),
        ("LANGUAGES",  ctx.get("snapshot", {}).get("languages", "—"),  False),
        ("PERIOD",     ctx.get("date_str", datetime.now().strftime("%B %Y")), False),
    ]
    for i, (lbl, val, hl) in enumerate(rows):
        ry = Inches(4.10 + i * 0.50)
        add_text(s, Inches(0.40), ry, Inches(0.95), Inches(0.26),
                 lbl, size=8, bold=True, color=GOLD_LT, char_spacing=2)
        if hl:
            add_rect(s, Inches(1.40), ry, Inches(2.05), Inches(0.30), GOLD_LT)
            add_text(s, Inches(1.45), ry + Inches(0.02),
                     Inches(2.00), Inches(0.26),
                     val, size=8, bold=True, color=NAVY_DK,
                     anchor=MSO_ANCHOR.MIDDLE)
        else:
            add_text(s, Inches(1.40), ry, Inches(2.05), Inches(0.30),
                     val, size=10, color=WHITE)

    # ─── RIGHT SIDE ───
    RX = LW + Inches(0.20)
    RW = SLIDE_W - LW - Inches(0.40)

    # National Snapshot strip title
    add_text(s, RX, Inches(0.30), RW, Inches(0.26),
             "NATIONAL SNAPSHOT  ·  " + ctx.get("date_str", datetime.now().strftime("%B %Y")),
             size=11, bold=True, color=GOLD, char_spacing=2)
    add_line(s, RX, Inches(0.58), RW, Emu(7000), GOLD, dashed=True)

    # 3-col hero snapshot
    hero_y = Inches(0.75)
    hero_h = Inches(2.30)
    col_w  = (RW - Inches(0.20)) / 3
    hero_cols = ctx.get("hero_snapshot", [])
    for i in range(3):
        cx = RX + col_w * i + (Inches(0.10) if i > 0 else Inches(0))
        cw = col_w - (Inches(0.10) if i < 2 else Inches(0))
        col_data = hero_cols[i] if i < len(hero_cols) else {}
        add_rect(s, cx, hero_y, cw, hero_h, WHITE, line=BORDER)
        add_text(s, cx + Inches(0.18), hero_y + Inches(0.10),
                 cw - Inches(0.36), Inches(0.24),
                 col_data.get("section", "—").upper(),
                 size=10, bold=True, color=GOLD, char_spacing=2)
        # Hero tile (top)
        tile_y = hero_y + Inches(0.42)
        add_text(s, cx + Inches(0.18), tile_y, cw - Inches(0.36), Inches(0.50),
                 col_data.get("hero_value", "—"),
                 size=28, bold=True, color=NAVY,
                 anchor=MSO_ANCHOR.BOTTOM)
        add_text(s, cx + Inches(0.18), tile_y + Inches(0.55),
                 cw - Inches(0.36), Inches(0.30),
                 col_data.get("hero_label", ""),
                 size=8.5, color=FG2)
        add_text(s, cx + Inches(0.18), tile_y + Inches(0.88),
                 cw - Inches(0.36), Inches(0.18),
                 col_data.get("hero_source", ""),
                 size=7, color=FG3, italic=True)
        # Sub-grid (2 cells)
        sg_y = tile_y + Inches(1.10)
        sub_cells = col_data.get("sub_cells", [])
        for j in range(2):
            if j >= len(sub_cells):
                continue
            sub = sub_cells[j]
            sub_x = cx + Inches(0.18) + (cw - Inches(0.36)) / 2 * j
            sub_w = (cw - Inches(0.40)) / 2
            add_text(s, sub_x, sg_y, sub_w, Inches(0.25),
                     sub.get("value", "—"), size=13, bold=True, color=NAVY)
            add_text(s, sub_x, sg_y + Inches(0.27), sub_w, Inches(0.20),
                     sub.get("label", ""), size=7.5, color=FG2)
            add_text(s, sub_x, sg_y + Inches(0.45), sub_w, Inches(0.16),
                     sub.get("source", ""), size=6.5, color=FG3, italic=True)

    # Geographic Context strip title
    geo_y = hero_y + hero_h + Inches(0.20)
    add_text(s, RX, geo_y, RW, Inches(0.26),
             "GEOGRAPHIC CONTEXT  ·  SECTOR STATUS",
             size=11, bold=True, color=GOLD, char_spacing=2)
    add_line(s, RX, geo_y + Inches(0.28), RW, Emu(7000), GOLD, dashed=True)

    # Lower row: map + (kpi pair + sector status)
    lower_y = geo_y + Inches(0.40)
    lower_h = Inches(3.20)
    map_w   = RW * 0.45
    side_w  = RW - map_w - Inches(0.15)
    # Map (image bytes injected by the renderer; falls back to placeholder if absent)
    add_map_placeholder(s, RX, lower_y, map_w, lower_h,
                        country=ctx["country"], lat=ctx.get("lat"), lng=ctx.get("lng"),
                        image_bytes=_ctx_map_image(ctx))
    # KPI pair
    side_x = RX + map_w + Inches(0.15)
    kpi_pair = ctx.get("kpi_pair", [])
    kpi_h = Inches(0.92)
    kpi_col_w = (side_w - Inches(0.10)) / 2
    for i in range(2):
        kx = side_x + (kpi_col_w + Inches(0.10)) * i
        k = kpi_pair[i] if i < len(kpi_pair) else {}
        add_rect(s, kx, lower_y, kpi_col_w, kpi_h, WHITE, line=BORDER)
        add_text(s, kx + Inches(0.15), lower_y + Inches(0.05),
                 kpi_col_w - Inches(0.30), Inches(0.22),
                 k.get("label", "").upper(), size=8, bold=True, color=GOLD, char_spacing=2)
        add_text(s, kx + Inches(0.15), lower_y + Inches(0.25),
                 kpi_col_w - Inches(0.30), Inches(0.42),
                 k.get("value", "—"), size=22, bold=True, color=NAVY)
        add_text(s, kx + Inches(0.15), lower_y + Inches(0.65),
                 kpi_col_w - Inches(0.30), Inches(0.20),
                 k.get("sublabel", ""), size=8.5, color=FG2)
    # Sector status panel
    sp_y = lower_y + kpi_h + Inches(0.10)
    sp_h = lower_h - kpi_h - Inches(0.10)
    add_rect(s, side_x, sp_y, side_w, sp_h, WHITE, line=BORDER)
    add_text(s, side_x + Inches(0.15), sp_y + Inches(0.08),
             side_w - Inches(0.30), Inches(0.24),
             "SECTOR STATUS  ·  " + ctx.get("date_str", datetime.now().strftime("%B %Y")),
             size=9, bold=True, color=GOLD, char_spacing=2)
    sst = ctx.get("sector_status", [])
    rows_y = sp_y + Inches(0.38)
    rows_h = sp_h - Inches(0.50)
    if sst:
        rh = rows_h / len(sst)
        for i, row in enumerate(sst):
            ry = rows_y + rh * i
            add_text(s, side_x + Inches(0.15), ry,
                     side_w * 0.35, rh,
                     row.get("name", ""), size=10, bold=True, color=NAVY,
                     anchor=MSO_ANCHOR.MIDDLE)
            add_text(s, side_x + Inches(0.15), ry + rh / 2,
                     side_w * 0.55 - Inches(0.20), rh / 2,
                     row.get("summary", ""), size=8, color=FG2, italic=True,
                     anchor=MSO_ANCHOR.TOP)
            # Chip
            status = (row.get("status") or "developing").lower()
            chip_color = STATUS_COLORS.get(status, GOLD)
            chip_w = Inches(0.90)
            chip_x = side_x + side_w - chip_w - Inches(0.15)
            chip_y = ry + rh / 2 - Inches(0.12)
            add_rect(s, chip_x, chip_y, chip_w, Inches(0.24), chip_color)
            add_text(s, chip_x, chip_y, chip_w, Inches(0.24),
                     status.title(), size=8, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Footer
    add_text(s, RX, SLIDE_H - Inches(0.40), RW, Inches(0.24),
             ctx.get("cover_footer",
                     "Sources: UN DESA · IMF · World Bank · OCHA · UNICEF · WHO"),
             size=8, color=FG3, italic=True)


def _sector_slide(prs, ctx, *, crumb, topic, statement, viz_renderer,
                  viz_title, stats, insights, interventions, sources_footer):
    """Standard sector slide (slides 2, 3, 4, 5, 6, 7, 8, 9, 10)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_header(s, crumb, topic, statement)
    # Left column — viz
    left_x = Inches(0.55); left_y = Inches(1.75)
    left_w = Inches(6.5);  left_h = Inches(5.10)
    add_section_label(s, left_x, left_y, left_w, viz_title)
    viz_renderer(s, left_x, left_y + Inches(0.42),
                 left_w, left_h - Inches(0.50))
    # Right column — stats + insights + interventions
    right_x = Inches(7.25); right_w = Inches(5.55)
    right_y = Inches(1.75)
    add_section_label(s, right_x, right_y, right_w, "Headline Indicators")
    stats_y = right_y + Inches(0.42)
    stats_h = Inches(0.72) * len(stats[:4])
    add_stats_block(s, right_x, stats_y, right_w, stats)
    insights_y = stats_y + stats_h + Inches(0.18)
    insights_h = Inches(1.45)
    add_insights_block(s, right_x, insights_y, right_w, insights_h, insights)
    intv_y = insights_y + insights_h + Inches(0.15)
    intv_h = Inches(6.90) - intv_y
    if intv_h > Inches(0.7):
        add_intv_block(s, right_x, intv_y, right_w, intv_h, interventions)
    add_footer(s, sources_footer)
    return s


def slide_2_at_a_glance(prs, ctx):
    sec = ctx.get("sectors", {}).get("at_a_glance", {})
    def viz(s, x, y, w, h):
        add_map_placeholder(s, x, y, w, h,
                            country=ctx["country"], lat=ctx.get("lat"),
                            image_bytes=_ctx_map_image(ctx))
    _sector_slide(prs, ctx,
        crumb="02 · NATIONAL OVERVIEW",
        topic="At a Glance",
        statement=sec.get("statement", "Country overview"),
        viz_renderer=viz,
        viz_title=sec.get("chart_title", "Regional Population Distribution"),
        stats=sec.get("stats", []),
        insights=sec.get("insights", []),
        interventions=sec.get("interventions", []),
        sources_footer=sec.get("sources", "Sources: UN DESA · World Bank"))


def slide_3_economy(prs, ctx):
    sec = ctx.get("sectors", {}).get("economy", {})
    chart = sec.get("chart", {})
    def viz(s, x, y, w, h):
        add_bar_line_chart(s, x, y, w, h,
            title=chart.get("title", "GDP vs Inflation"),
            x_labels=chart.get("x_labels", ["FY20","FY21","FY22","FY23","FY24"]),
            bar_vals=chart.get("bar_values", [6.1,6.3,6.4,6.6,7.3]),
            line_vals=chart.get("line_values", [14,16,24,28,26.6]),
            y_max=chart.get("y_max", 40),
            bar_label=chart.get("bar_label", "GDP growth"),
            line_label=chart.get("line_label", "CPI inflation"))
    _sector_slide(prs, ctx,
        crumb="03 · ECONOMY & DEBT",
        topic="Economy",
        statement=sec.get("statement", ""),
        viz_renderer=viz,
        viz_title=chart.get("title", "GDP Growth vs. Inflation"),
        stats=sec.get("stats", []),
        insights=sec.get("insights", []),
        interventions=sec.get("interventions", []),
        sources_footer=sec.get("sources", "Sources: AfDB · World Bank · IMF · NBE"))


def slide_4_health(prs, ctx):
    sec = ctx.get("sectors", {}).get("health", {})
    def viz(s, x, y, w, h):
        add_map_placeholder(s, x, y, w, h, country=ctx["country"],
            legend=[("Low coverage", NEG), ("Mid", GOLD), ("High", POS)],
            image_bytes=_ctx_map_image(ctx))
    _sector_slide(prs, ctx,
        crumb="04 · HEALTH",
        topic="Health",
        statement=sec.get("statement", ""),
        viz_renderer=viz,
        viz_title=sec.get("chart_title", "Full Immunization Coverage by Region"),
        stats=sec.get("stats", []),
        insights=sec.get("insights", []),
        interventions=sec.get("interventions", []),
        sources_footer=sec.get("sources", "Sources: WHO · UNICEF · UN IGME"))


def slide_5_education(prs, ctx):
    sec = ctx.get("sectors", {}).get("education", {})
    chart = sec.get("chart", {})
    def viz(s, x, y, w, h):
        add_grouped_bars(s, x, y, w, h,
            title=chart.get("title", "Enrollment by Level & Gender"),
            x_labels=chart.get("x_labels", ["Primary","Secondary","Tertiary"]),
            series_a=chart.get("series_a", [95, 32, 8]),
            series_b=chart.get("series_b", [88, 28, 6]),
            a_label=chart.get("a_label", "Boys"),
            b_label=chart.get("b_label", "Girls"),
            y_max=chart.get("y_max", 100))
    _sector_slide(prs, ctx,
        crumb="05 · EDUCATION",
        topic="Education",
        statement=sec.get("statement", ""),
        viz_renderer=viz,
        viz_title=chart.get("title", "Enrollment by Level & Gender"),
        stats=sec.get("stats", []),
        insights=sec.get("insights", []),
        interventions=sec.get("interventions", []),
        sources_footer=sec.get("sources", "Sources: UNESCO UIS · Ministry of Education"))


def slide_6_nutrition(prs, ctx):
    sec = ctx.get("sectors", {}).get("nutrition", {})
    funnel = sec.get("funnel", [])
    def viz(s, x, y, w, h):
        add_funnel(s, x, y, w, h,
            title=sec.get("chart_title", "Population Cascade — Risk to Acute Need"),
            rows=funnel)
    _sector_slide(prs, ctx,
        crumb="06 · FOOD SECURITY & NUTRITION",
        topic="Nutrition Cascade",
        statement=sec.get("statement", ""),
        viz_renderer=viz,
        viz_title=sec.get("chart_title", "Population Cascade"),
        stats=sec.get("stats", []),
        insights=sec.get("insights", []),
        interventions=sec.get("interventions", []),
        sources_footer=sec.get("sources", "Sources: IPC · FEWS NET · WFP"))


def slide_7_agriculture(prs, ctx):
    sec = ctx.get("sectors", {}).get("agriculture", {})
    def viz(s, x, y, w, h):
        add_map_placeholder(s, x, y, w, h, country=ctx["country"],
            legend=[("<20%", POS), ("20–35%", GOLD), ("35%+", NEG)],
            image_bytes=_ctx_map_image(ctx))
    _sector_slide(prs, ctx,
        crumb="07 · AGRICULTURE & RURAL LIVELIHOODS",
        topic="Agriculture",
        statement=sec.get("statement", ""),
        viz_renderer=viz,
        viz_title=sec.get("chart_title", "Stunting Prevalence by Region"),
        stats=sec.get("stats", []),
        insights=sec.get("insights", []),
        interventions=sec.get("interventions", []),
        sources_footer=sec.get("sources", "Sources: FAOSTAT · CSA · MoA"))


def slide_8_infrastructure(prs, ctx):
    sec = ctx.get("sectors", {}).get("infrastructure", {})
    def viz(s, x, y, w, h):
        add_compare_rows(s, x, y, w, h,
            title=sec.get("chart_title", "Service Access — Urban vs. Rural"),
            rows=sec.get("compare_rows", [
                {"label":"Electricity",    "urban":95, "rural":15},
                {"label":"Safe water",     "urban":85, "rural":50},
                {"label":"Improved sanit.","urban":35, "rural":12},
                {"label":"Internet",       "urban":45, "rural":8},
            ]))
    _sector_slide(prs, ctx,
        crumb="08 · INFRASTRUCTURE & CONNECTIVITY",
        topic="Infrastructure",
        statement=sec.get("statement", ""),
        viz_renderer=viz,
        viz_title=sec.get("chart_title", "Service Access — Urban vs. Rural"),
        stats=sec.get("stats", []),
        insights=sec.get("insights", []),
        interventions=sec.get("interventions", []),
        sources_footer=sec.get("sources", "Sources: World Bank · ITU"))


def slide_9_climate(prs, ctx):
    sec = ctx.get("sectors", {}).get("climate", {})
    chart = sec.get("chart", {})
    def viz(s, x, y, w, h):
        add_line_chart(s, x, y, w, h,
            title=chart.get("title", "Forest Cover Decline & Climate Risk"),
            x_labels=chart.get("x_labels", ["2000","2005","2010","2015","2020","2024"]),
            series=chart.get("series", [
                {"label":"Forest cover (%)",       "values":[40,35,28,22,17,15], "color":DEEP_GREEN},
                {"label":"Climate risk (ND-GAIN)", "values":[20,25,33,45,58,68], "color":NEG, "dashed":True},
            ]),
            y_max=chart.get("y_max", 100))
    _sector_slide(prs, ctx,
        crumb="09 · ENVIRONMENT & CLIMATE",
        topic="Climate",
        statement=sec.get("statement", ""),
        viz_renderer=viz,
        viz_title=chart.get("title", "Forest Cover Decline & Climate Risk"),
        stats=sec.get("stats", []),
        insights=sec.get("insights", []),
        interventions=sec.get("interventions", []),
        sources_footer=sec.get("sources", "Sources: World Bank · ND-GAIN · IPCC"))


def slide_10_humanitarian(prs, ctx):
    sec = ctx.get("sectors", {}).get("humanitarian", {})
    def viz(s, x, y, w, h):
        add_map_placeholder(s, x, y, w, h, country=ctx["country"],
            legend=[("<100K", BLUE), ("100K–500K", GOLD), ("500K+", NEG)],
            image_bytes=_ctx_map_image(ctx))
    _sector_slide(prs, ctx,
        crumb="10 · HUMANITARIAN",
        topic="Humanitarian",
        statement=sec.get("statement", ""),
        viz_renderer=viz,
        viz_title=sec.get("chart_title", "IDP Concentration by Region"),
        stats=sec.get("stats", []),
        insights=sec.get("insights", []),
        interventions=sec.get("interventions", []),
        sources_footer=sec.get("sources", "Sources: IOM DTM · UNHCR · OCHA"))


def slide_11_bigstat(prs, ctx):
    """Full-bleed navy slide — funding-gap headline."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    bs = ctx.get("bigstat", {})
    # Top crumb
    add_text(s, Inches(0.6), Inches(0.45), SLIDE_W - Inches(1.2), Inches(0.30),
             bs.get("crumb", "11 · FUNDING GAP").upper(),
             size=11, bold=True, color=GOLD_LT, char_spacing=3)
    add_line(s, Inches(0.6), Inches(0.80), Inches(2.0), Emu(8000), GOLD)
    # Body — left big number, right statement+quote
    body_y = Inches(1.40); body_h = Inches(4.50)
    # Left
    add_text(s, Inches(0.6), body_y, Inches(6.0), Inches(2.40),
             bs.get("hero_value", "~40%"),
             size=140, bold=True, color=GOLD)
    add_text(s, Inches(0.6), body_y + Inches(2.45), Inches(6.0), Inches(0.50),
             bs.get("hero_label", ""), size=14, bold=True,
             color=GOLD_LT, char_spacing=2)
    # Right
    rx = Inches(7.20)
    add_text(s, rx, body_y, SLIDE_W - rx - Inches(0.6), Inches(2.20),
             bs.get("statement", ""),
             size=22, color=WHITE)
    # Quote
    q_y = body_y + Inches(2.40)
    add_rect(s, rx, q_y, Inches(0.06), Inches(1.30), GOLD)
    add_text(s, rx + Inches(0.20), q_y,
             SLIDE_W - rx - Inches(0.6) - Inches(0.20), Inches(1.10),
             bs.get("quote", ""), size=14, italic=True, color=GOLD_LT)
    add_text(s, rx + Inches(0.20), q_y + Inches(1.05),
             SLIDE_W - rx - Inches(0.6) - Inches(0.20), Inches(0.30),
             bs.get("quote_source", ""), size=10, color=GOLD_LT, italic=True)
    # Bottom stat strip
    strip_y = Inches(6.00)
    strip_h = Inches(1.00)
    add_rect(s, 0, strip_y, SLIDE_W, strip_h, NAVY_DK)
    cells = bs.get("strip", [])
    if cells:
        cw = SLIDE_W / len(cells)
        for i, c in enumerate(cells):
            cx = cw * i
            add_text(s, cx, strip_y + Inches(0.15), cw, Inches(0.50),
                     c.get("value", ""), size=30, bold=True, color=GOLD,
                     align=PP_ALIGN.CENTER)
            add_text(s, cx, strip_y + Inches(0.65), cw, Inches(0.28),
                     c.get("label", "").upper(), size=10, color=GOLD_LT,
                     align=PP_ALIGN.CENTER, char_spacing=2)


def slide_12_priorities(prs, ctx):
    """8-card recommendations grid (4×2)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, BG)
    pri = ctx.get("priorities_block", {})
    # Header
    add_text(s, Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.30),
             pri.get("crumb", "12 · STRATEGIC RECOMMENDATIONS").upper(),
             size=11, bold=True, color=BLUE, char_spacing=2.5)
    add_runs(s, Inches(0.55), Inches(0.65), Inches(12.2), Inches(0.85),
             [{"text": pri.get("topic", "Priorities") + " ",
               "size": 30, "bold": True, "color": NAVY},
              {"text": "|  ", "size": 30, "color": GOLD},
              {"text": pri.get("statement", ""), "size": 30, "bold": True, "color": GOLD}],
             anchor=MSO_ANCHOR.MIDDLE)
    # Grid
    cards = (pri.get("cards") or [])[:8]
    grid_x = Inches(0.55); grid_y = Inches(1.70)
    grid_w = Inches(12.20); grid_h = Inches(4.70)
    cols = 4; rows = 2
    gap = Inches(0.18)
    cell_w = (grid_w - gap * (cols - 1)) / cols
    cell_h = (grid_h - gap * (rows - 1)) / rows
    TAG_COLORS = {
        "high":         (NEG,  GOLD_PALE),
        "medium":       (NEUTRAL, GOLD_PALE),
        "foundational": (BLUE, BLUE_PALE),
    }
    for i, card in enumerate(cards):
        r = i // cols; c = i % cols
        cx = grid_x + (cell_w + gap) * c
        cy = grid_y + (cell_h + gap) * r
        dark = card.get("dark", False)
        bg_color = NAVY if dark else WHITE
        accent   = GOLD if dark else BLUE
        add_rect(s, cx, cy, cell_w, cell_h, bg_color, line=BORDER)
        add_rect(s, cx, cy, cell_w, Inches(0.08), accent)
        add_text(s, cx + Inches(0.20), cy + Inches(0.20),
                 cell_w - Inches(0.40), Inches(0.24),
                 card.get("num", f"{i+1:02d} · {card.get('sector','')}"),
                 size=9.5, bold=True, color=GOLD_LT if dark else GOLD,
                 char_spacing=2)
        add_text(s, cx + Inches(0.20), cy + Inches(0.50),
                 cell_w - Inches(0.40), Inches(0.70),
                 card.get("title", ""),
                 size=14, bold=True, color=WHITE if dark else NAVY)
        add_text(s, cx + Inches(0.20), cy + Inches(1.20),
                 cell_w - Inches(0.40), cell_h - Inches(1.60),
                 card.get("desc", ""),
                 size=10, color=GOLD_LT if dark else FG2)
        # Priority tag
        tag = (card.get("priority") or "high").lower()
        tag_color, tag_bg = TAG_COLORS.get(tag, TAG_COLORS["high"])
        tag_label = {"high":"● HIGH", "medium":"● MEDIUM",
                     "foundational":"● FOUNDATIONAL"}.get(tag, "● HIGH")
        tag_y = cy + cell_h - Inches(0.40)
        add_rect(s, cx + Inches(0.20), tag_y, Inches(1.40), Inches(0.26),
                 tag_bg)
        add_text(s, cx + Inches(0.20), tag_y, Inches(1.40), Inches(0.26),
                 tag_label, size=9, bold=True, color=tag_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, char_spacing=1)
    # Footer
    add_line(s, Inches(0.55), Inches(6.55), Inches(12.2), Emu(4000), BORDER)
    add_text(s, Inches(0.55), Inches(6.65), Inches(12.2), Inches(0.40),
             pri.get("footer",
                 "Priority levels reflect urgency × systemic impact, not preference."),
             size=10, color=FG2, italic=True)


def slide_13_closing(prs, ctx):
    """Closing — full navy."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    c = ctx.get("closing", {})
    # Gold bar
    add_rect(s, SLIDE_W / 2 - Inches(0.40), Inches(2.00),
             Inches(0.80), Inches(0.05), GOLD)
    # Country name
    add_text(s, Inches(1.0), Inches(2.20), SLIDE_W - Inches(2.0), Inches(1.60),
             ctx["country"], size=66, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Subtitle
    add_text(s, Inches(1.0), Inches(4.00), SLIDE_W - Inches(2.0), Inches(0.50),
             c.get("subtitle", "Country Overview · Sector Assessment & Priorities"),
             size=18, color=GOLD_LT, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    # Meta
    meta = c.get("meta", [])
    add_text(s, Inches(1.0), Inches(4.70), SLIDE_W - Inches(2.0), Inches(0.30),
             "  ·  ".join(meta) if meta else
             ctx.get("date_str", datetime.now().strftime("%B %Y")),
             size=10, color=GOLD_LT, align=PP_ALIGN.CENTER, char_spacing=3)
    # Sources block
    add_text(s, Inches(1.5), Inches(5.50), SLIDE_W - Inches(3.0), Inches(1.50),
             c.get("sources",
                 "Sources: World Bank · IMF · OCHA · WHO · UNICEF · UNESCO UIS · UNHCR · IOM DTM"),
             size=10, color=GOLD_LT, italic=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# ENTRY POINT
# ═══════════════════════════════════════════════════════════════════════════════

def build(context: dict, output=None):
    """Build the 13-slide country brief."""
    ctx = _safe_copy(context)
    ctx.setdefault("date_str", datetime.now().strftime("%B %Y"))
    ctx.setdefault("sectors", {})

    # ─── Auto-render map if not already present (all 3 paths get maps) ──
    # Caller may pre-populate ctx['map']['image_bytes'] (e.g. telegram_webhook
    # does this so it can report X-Map-Status); otherwise we attempt the
    # renderer here. Silent fallback if the module is missing.
    if not (ctx.get("map") or {}).get("image_bytes") and ctx.get("country"):
        try:
            from country_map_renderer import render as _render_map
            _png = _render_map(
                country=ctx["country"],
                iso2=ctx.get("iso2", ""),
                lat=ctx.get("lat"),
                lng=ctx.get("lng"),
                type=(ctx.get("map") or {}).get("type", "reference"),
                indicators=ctx.get("indicators"),
                subnational_indicators=ctx.get("subnational_indicators"),
            )
            if _png:
                ctx.setdefault("map", {})
                ctx["map"]["image_bytes"] = _png
        except Exception:
            pass   # renderer unavailable → builder falls back to placeholder

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_1_cover(prs, ctx)
    slide_2_at_a_glance(prs, ctx)
    slide_3_economy(prs, ctx)
    slide_4_health(prs, ctx)
    slide_5_education(prs, ctx)
    slide_6_nutrition(prs, ctx)
    slide_7_agriculture(prs, ctx)
    slide_8_infrastructure(prs, ctx)
    slide_9_climate(prs, ctx)
    slide_10_humanitarian(prs, ctx)
    slide_11_bigstat(prs, ctx)
    slide_12_priorities(prs, ctx)
    slide_13_closing(prs, ctx)

    if output is None:
        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()
    if hasattr(output, "write"):
        prs.save(output)
        return None
    prs.save(output)
    return None


CONTEXT_SCHEMA = """
context = {
  "country":   "Ethiopia",
  "iso2":      "et",
  "lat": 9.145, "lng": 40.4897,
  "date_str":  "May 2026",
  "subtitle":  "Federal Democratic Republic of Ethiopia · Horn of Africa · 130M people",

  "snapshot": {
      "capital": "Addis Ababa · 5.2M",
      "currency":"Ethiopian Birr (ETB)",
      "languages":"Amharic · Oromo · Tigrinya · Somali",
      "dac_class":"LEAST DEVELOPED · LOW INCOME",
      "gni_cap": "$1,020 (2024)",
  },

  "hero_snapshot": [
      {"section":"General",
       "hero_value":"$1.1B","hero_label":"UAE contribution · aid & deposits 2018–24",
       "hero_source":"UAE MoFA · 2024",
       "sub_cells":[
           {"value":"0.50","label":"HDI · rank 175","source":"UNDP · 2024"},
           {"value":"~250K","label":"Diaspora in UAE","source":"IOM · MoLSA"},
       ]},
      ...
  ],

  "kpi_pair": [
      {"label":"Population","value":"130M","sublabel":"2nd in Africa  ·  UN DESA 2024"},
      {"label":"Land Area","value":"1.1M km²","sublabel":"12th in Africa  ·  landlocked"},
  ],

  "sector_status": [   # 7 rows on cover
      {"name":"Health","status":"severe","summary":"High maternal/child mortality"},
      ...
  ],

  "cover_footer": "Sources: UN DESA · IMF · World Bank · AfDB · UNDP · OCHA",

  "sectors": {
      "at_a_glance": {statement, chart_title, stats, insights, interventions, sources},
      "economy":     {statement, chart{title,x_labels,bar_values,line_values,y_max},
                      stats, insights, interventions, sources},
      "health":      {...},
      "education":   {..., chart{title,x_labels,series_a,series_b,a_label,b_label,y_max}},
      "nutrition":   {..., funnel:[{pct,percent_text,label,stage,sublabel,width_ratio,color}]},
      "agriculture": {...},
      "infrastructure":{..., compare_rows:[{label,urban,rural}]},
      "climate":     {..., chart{title,x_labels,series:[{label,values,color,dashed?}],y_max}},
      "humanitarian":{...},
  },

  "bigstat": {
      "crumb":"11 · FUNDING GAP · 2024 HRP",
      "hero_value":"~40%",  "hero_label":"OF $2.4B HRP FUNDED IN 2024",
      "statement":"$1.4B unfunded means food, shelter ...",
      "quote":"...",  "quote_source":"— OCHA Ethiopia, 2024",
      "strip":[{value,label}, ...]
  },

  "priorities_block": {
      "crumb":"12 · STRATEGIC RECOMMENDATIONS",
      "topic":"Priorities",
      "statement":"Eight cross-sector interventions, sequenced by urgency",
      "cards":[
          {"num":"01 · Economy","title":"...","desc":"...","priority":"high","dark":False},
          ...
      ],
      "footer":"..."
  },

  "closing": {
      "subtitle":"Country Overview · Sector Assessment & Intervention Priorities",
      "meta":["May 2026","Country Analysis Team"],
      "sources":"Sources: World Bank · ..."
  },
}
"""
